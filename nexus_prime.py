#!/usr/bin/env python3
# ╔══════════════════════════════════════════════════════════════════════════════╗
# ║                    N E X U S   P R I M E  v9.0 — WEB                       ║
# ║          Asistente Academico Universitario de Elite                         ║
# ║                                                                             ║
# ║  Motor Principal : Groq  (llama-3.3-70b-versatile)                         ║
# ║  Motor Respaldo  : Gemini 2.0 Flash                                         ║
# ║  Motor Extra     : OpenRouter (sin saturacion)                              ║
# ║  Imagenes        : Pollinations AI (gratis, sin clave)                      ║
# ║  Interfaz        : HTML moderna servida via Flask                           ║
# ║                                                                             ║
# ║  EJECUCION:                                                                 ║
# ║    python nexus_prime.py                                                    ║
# ║    Luego abre: http://localhost:8080                                        ║
# ╚══════════════════════════════════════════════════════════════════════════════╝

from __future__ import annotations

# ══════════════════════════════════════════════════════════════════════════════
#  AUTO-INSTALACION DE DEPENDENCIAS
# ══════════════════════════════════════════════════════════════════════════════
import subprocess, sys as _sys

_PAQUETES = [
    "flask",
    "flask-cors",
    "httpx==0.27.2",
    "groq",
    "google-genai",
    "pytz",
    "PyMuPDF",
    "python-docx",
    "openpyxl",
    "python-pptx",
    "Pillow",
    "beautifulsoup4",
    "lxml",
]


def _instalar_paquetes() -> None:
    print("🔧 Verificando dependencias...")
    try:
        subprocess.check_call(
            [_sys.executable, "-m", "pip", "install", "--quiet", "--upgrade"]
            + _PAQUETES,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        print("✅ Dependencias listas.")
    except subprocess.CalledProcessError as e:
        print(f"⚠️  Error al instalar dependencias: {e}")


_instalar_paquetes()

# ══════════════════════════════════════════════════════════════════════════════
import os
import re
import json
import logging
import asyncio
import base64
import urllib.parse
import threading
import datetime
from io import BytesIO
from typing import Optional, List, Dict, Any, Tuple
from collections import defaultdict
import xml.etree.ElementTree as ET

import httpx

try:
    import pytz as _pytz
    _TZ = _pytz.timezone("America/Lima")
    def _now_lima() -> datetime.datetime:
        return datetime.datetime.now(_TZ)
except ImportError:
    def _now_lima() -> datetime.datetime:
        return datetime.datetime.now(datetime.timezone.utc) - datetime.timedelta(hours=5)

# ══════════════════════════════════════════════════════════════════════════════
#  LOGGING
# ══════════════════════════════════════════════════════════════════════════════
logging.basicConfig(
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
    level=logging.INFO,
)
logger = logging.getLogger(__name__)

# ══════════════════════════════════════════════════════════════════════════════
#  LLAVES API — Configura aqui o como variables de entorno
# ══════════════════════════════════════════════════════════════════════════════
GROQ_API_KEY: str = os.environ.get("GROQ_API_KEY", "")
GEMINI_API_KEY: str = os.environ.get("GEMINI_API_KEY", "")
OPENROUTER_API_KEY: str = os.environ.get("OPENROUTER_API_KEY", "")

# ══════════════════════════════════════════════════════════════════════════════
#  CONSTANTES
# ══════════════════════════════════════════════════════════════════════════════
BOT_VERSION = "9.0"
BOT_NAME = "Nexus"
MAX_HISTORY = 30
GROQ_MODEL = "llama-3.3-70b-versatile"
GEMINI_MODEL = "gemini-2.0-flash"
OPENROUTER_MODEL = "mistralai/mistral-7b-instruct:free"
MAX_DOC_CHARS = 12000
PORT = int(os.environ.get("PORT", 8080))

# ══════════════════════════════════════════════════════════════════════════════
#  PROMPT DEL SISTEMA
# ══════════════════════════════════════════════════════════════════════════════
SYSTEM_PROMPT: str = (
    f"Eres {BOT_NAME}, un asistente academico universitario de elite. "
    "Tu personalidad es rigurosa, empatica, clara y pedagogica. "
    "Hablas con un tono profesional pero cercano, como un tutor de confianza.\n\n"
    "REGLAS:\n"
    "1. Responde SIEMPRE en espanol.\n"
    "2. Usa lenguaje universitario formal pero comprensible.\n"
    "3. Cita en formato APA 7 cuando menciones fuentes academicas.\n"
    "4. JAMAS inventes referencias, DOIs ni autores.\n"
    "5. Cuando des enlaces, asegurate de que sean de revistas cientificas indexadas.\n"
    "6. Prefiere fuentes de 2019 a 2026.\n"
    "7. Da respuestas completas, estructuradas y utiles.\n"
    "8. Puedes conversar libremente y ayudar con cualquier tema.\n"
    "9. Usa formato Markdown para estructurar las respuestas (##, **negrita**, listas).\n"
)

# ══════════════════════════════════════════════════════════════════════════════
#  ESTADO GLOBAL (por sesion web via session_id)
# ══════════════════════════════════════════════════════════════════════════════
conversation_history: Dict[str, List[Dict[str, str]]] = {}
user_mode: Dict[str, str] = {}
user_model: Dict[str, str] = {}
user_stats: Dict[str, int] = {}
user_tasks: Dict[str, List[Dict[str, Any]]] = defaultdict(list)
user_goals: Dict[str, str] = {}
doc_context: Dict[str, Dict[str, str]] = {}
paciente_sessions: Dict[str, Dict[str, Any]] = {}

# ══════════════════════════════════════════════════════════════════════════════
#  CLIENTES IA
# ══════════════════════════════════════════════════════════════════════════════
groq_client: Any = None
gemini_client: Any = None
gemini_types: Any = None


def _init_clients() -> None:
    global groq_client, gemini_client, gemini_types
    if GROQ_API_KEY:
        try:
            from groq import Groq as _Groq
            groq_client = _Groq(api_key=GROQ_API_KEY)
            logger.info("✅ Groq inicializado")
        except Exception as exc:
            logger.error("Error Groq: %s", exc)

    if GEMINI_API_KEY:
        try:
            from google import genai as _genai
            from google.genai import types as _gtypes
            gemini_client = _genai.Client(api_key=GEMINI_API_KEY)
            gemini_types = _gtypes
            logger.info("✅ Gemini inicializado")
        except Exception as exc:
            logger.error("Error Gemini: %s", exc)

    if OPENROUTER_API_KEY:
        logger.info("✅ OpenRouter listo")


_init_clients()

# ══════════════════════════════════════════════════════════════════════════════
#  MODOS
# ══════════════════════════════════════════════════════════════════════════════
MODE_SYSTEMS: Dict[str, str] = {
    "libre": (
        f"Eres {BOT_NAME}, asistente academico versatil. "
        "Responde en espanol de forma clara, util y amigable. "
        "Usa Markdown para estructurar tus respuestas. "
        "Cita en APA 7 cuando uses fuentes academicas."
    ),
    "academico": (
        f"Eres {BOT_NAME} en modo academico estricto. "
        "Lenguaje universitario formal siempre. Todas las respuestas incluyen citas APA 7. "
        "Estructura tus respuestas con titulos y subtitulos claros usando Markdown."
    ),
    "codigo": (
        f"Eres {BOT_NAME} en modo programacion. Experto en todos los lenguajes. "
        "Codigo limpio, comentado y con buenas practicas. Usa bloques ```codigo```. En espanol."
    ),
    "creativo": (
        f"Eres {BOT_NAME} en modo creativo. Escritor, poeta, narrador de elite. "
        "Contenido original, imaginativo y literariamente rico. Siempre en espanol."
    ),
    "tutor": (
        f"Eres {BOT_NAME} en modo tutor. Explica como un profesor universitario paciente. "
        "Usa ejemplos reales, analogias y ejercicios practicos. En espanol."
    ),
    "clinico": (
        f"Eres {BOT_NAME} en modo PSICOLOGO CLINICO. "
        "Simulas un paciente con caracteristicas psicopatologicas especificas. "
        "Responde SIEMPRE como el paciente (primera persona). "
        "No rompas el personaje a menos que el usuario escriba /salirclinico."
    ),
}

# ══════════════════════════════════════════════════════════════════════════════
#  EXTRACCION DE TEXTO DE DOCUMENTOS
# ══════════════════════════════════════════════════════════════════════════════
def extract_text_from_file(file_bytes: bytes, mime_type: str, filename: str) -> Tuple[str, str]:
    fn = filename.lower()
    texto = ""
    tipo = "desconocido"

    try:
        if "pdf" in mime_type or fn.endswith(".pdf"):
            tipo = "PDF"
            try:
                import pymupdf as _fitz
            except ImportError:
                import fitz as _fitz
            pdf_doc = _fitz.open(stream=file_bytes, filetype="pdf")
            partes = []
            for i, page in enumerate(pdf_doc):
                if i >= 30:
                    partes.append(f"\n[... documento continúa, mostrando primeras 30 páginas ...]")
                    break
                partes.append(page.get_text())
            pdf_doc.close()
            texto = "\n".join(partes)

        elif "word" in mime_type or fn.endswith(".docx"):
            tipo = "Word (DOCX)"
            import docx as _python_docx
            word_doc = _python_docx.Document(BytesIO(file_bytes))
            parrafos = [p.text for p in word_doc.paragraphs if p.text and p.text.strip()]
            for table in word_doc.tables:
                for row in table.rows:
                    fila = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if fila:
                        parrafos.append(fila)
            texto = "\n".join(parrafos)

        elif "spreadsheet" in mime_type or fn.endswith((".xlsx", ".xls")):
            tipo = "Excel (XLSX)"
            import openpyxl as _openpyxl
            wb = _openpyxl.load_workbook(BytesIO(file_bytes), read_only=True, data_only=True)
            filas: List[str] = []
            for sheet_name in wb.sheetnames[:5]:
                ws = wb[sheet_name]
                filas.append(f"=== Hoja: {sheet_name} ===")
                for row in ws.iter_rows(max_row=200, values_only=True):
                    row_text = " | ".join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        filas.append(row_text)
            wb.close()
            texto = "\n".join(filas)

        elif "presentation" in mime_type or fn.endswith(".pptx"):
            tipo = "PowerPoint (PPTX)"
            from pptx import Presentation as _Presentation
            prs = _Presentation(BytesIO(file_bytes))
            diapositivas: List[str] = []
            for i, slide in enumerate(prs.slides):
                textos_slide: List[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = str(shape.text).strip()
                        if t:
                            textos_slide.append(t)
                if textos_slide:
                    diapositivas.append(f"--- Diapositiva {i + 1} ---\n" + "\n".join(textos_slide))
            texto = "\n\n".join(diapositivas)

        elif mime_type in ("text/plain", "text/csv", "text/markdown") or fn.endswith((".txt", ".csv", ".md")):
            tipo = "Texto plano"
            texto = file_bytes.decode("utf-8", errors="replace")

        else:
            tipo = "No soportado"
            texto = ""

    except Exception as exc:
        logger.error("Error extrayendo texto de %s: %s", filename, exc)
        tipo = "Error al leer"
        texto = ""

    if len(texto) > MAX_DOC_CHARS:
        texto = texto[:MAX_DOC_CHARS] + f"\n\n[... texto recortado a {MAX_DOC_CHARS} caracteres ...]"

    return texto.strip(), tipo


# ══════════════════════════════════════════════════════════════════════════════
#  MOTORES IA
# ══════════════════════════════════════════════════════════════════════════════
async def query_groq(sid: str, prompt: str, system: str) -> Optional[str]:
    if groq_client is None:
        return None
    hist = conversation_history.setdefault(sid, [])
    hist.append({"role": "user", "content": prompt})
    if len(hist) > MAX_HISTORY:
        conversation_history[sid] = hist[-MAX_HISTORY:]

    try:
        from groq.types.chat import (
            ChatCompletionSystemMessageParam,
            ChatCompletionUserMessageParam,
            ChatCompletionAssistantMessageParam,
        )
        messages_typed: List[Any] = [
            ChatCompletionSystemMessageParam(role="system", content=system)
        ]
        for m in conversation_history[sid]:
            if m["role"] == "user":
                messages_typed.append(ChatCompletionUserMessageParam(role="user", content=m["content"]))
            elif m["role"] == "assistant":
                messages_typed.append(ChatCompletionAssistantMessageParam(role="assistant", content=m["content"]))
    except ImportError:
        messages_typed = [{"role": "system", "content": system}] + conversation_history[sid]

    waits = [5, 15, 30]
    for attempt in range(3):
        try:
            def _call() -> str:
                resp = groq_client.chat.completions.create(
                    model=GROQ_MODEL,
                    messages=messages_typed,
                    max_tokens=4096,
                    temperature=0.7,
                )
                choices = resp.choices
                if not choices:
                    return ""
                m = choices[0].message
                return (m.content or "") if m else ""

            text: str = await asyncio.to_thread(_call)
            if text.strip():
                conversation_history[sid].append({"role": "assistant", "content": text})
                return text
            return None
        except Exception as exc:
            err = str(exc).lower()
            logger.warning("Groq intento %d/3: %s", attempt + 1, exc)
            if any(k in err for k in ["rate", "429", "too many", "quota"]):
                if attempt < 2:
                    await asyncio.sleep(waits[attempt])
                    continue
                return None
            if any(k in err for k in ["invalid", "401", "403", "unauthorized"]):
                return None
            if attempt < 2:
                await asyncio.sleep(waits[attempt])
    return None


async def query_gemini(sid: str, prompt: str, system: str) -> Optional[str]:
    if gemini_client is None or gemini_types is None:
        return None
    try:
        contents = [
            gemini_types.Content(role="user", parts=[gemini_types.Part(text=prompt)])
        ]
        config = gemini_types.GenerateContentConfig(
            system_instruction=system,
            temperature=0.7,
            max_output_tokens=4096,
        )

        def _call() -> str:
            resp = gemini_client.models.generate_content(
                model=GEMINI_MODEL, contents=contents, config=config,
            )
            result = ""
            if resp.candidates:
                for cand in resp.candidates:
                    if cand.content and cand.content.parts:
                        for part in cand.content.parts:
                            if hasattr(part, "text") and part.text:
                                result += part.text
            return result

        text: str = await asyncio.to_thread(_call)
        if text.strip():
            conversation_history.setdefault(sid, [])
            conversation_history[sid].append({"role": "user", "content": prompt})
            conversation_history[sid].append({"role": "assistant", "content": text})
            if len(conversation_history[sid]) > MAX_HISTORY:
                conversation_history[sid] = conversation_history[sid][-MAX_HISTORY:]
            return text.strip()
        return None
    except Exception as exc:
        logger.warning("Gemini error: %s", exc)
        return None


async def query_gemini_vision(image_bytes: bytes, mime_type: str, prompt: str) -> Optional[str]:
    if gemini_client is None or gemini_types is None:
        return None
    try:
        contents = [
            gemini_types.Content(
                role="user",
                parts=[
                    gemini_types.Part(
                        inline_data=gemini_types.Blob(mime_type=mime_type, data=image_bytes)
                    ),
                    gemini_types.Part(text=prompt),
                ],
            )
        ]
        config = gemini_types.GenerateContentConfig(
            system_instruction=SYSTEM_PROMPT, temperature=0.5, max_output_tokens=4096,
        )

        def _call() -> str:
            resp = gemini_client.models.generate_content(
                model=GEMINI_MODEL, contents=contents, config=config,
            )
            result = ""
            if resp.candidates:
                for cand in resp.candidates:
                    if cand.content and cand.content.parts:
                        for part in cand.content.parts:
                            if hasattr(part, "text") and part.text:
                                result += part.text
            return result

        text: str = await asyncio.to_thread(_call)
        return text.strip() if text.strip() else None
    except Exception as exc:
        logger.warning("Gemini Vision error: %s", exc)
        return None


async def query_openrouter(sid: str, prompt: str, system: str) -> Optional[str]:
    if not OPENROUTER_API_KEY:
        return None
    hist = conversation_history.setdefault(sid, [])
    hist.append({"role": "user", "content": prompt})
    if len(hist) > MAX_HISTORY:
        conversation_history[sid] = hist[-MAX_HISTORY:]

    messages = [{"role": "system", "content": system}] + conversation_history[sid]
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {OPENROUTER_API_KEY}",
                    "Content-Type": "application/json",
                    "HTTP-Referer": "https://nexus-academic.app",
                    "X-Title": "Nexus Academic Bot",
                },
                json={
                    "model": OPENROUTER_MODEL,
                    "messages": messages,
                    "max_tokens": 4096,
                    "temperature": 0.7,
                },
            )
            if resp.status_code == 200:
                data = resp.json()
                text = data["choices"][0]["message"]["content"]
                if text and text.strip():
                    conversation_history[sid].append({"role": "assistant", "content": text})
                    return text.strip()
            else:
                logger.warning("OpenRouter status %d: %s", resp.status_code, resp.text[:200])
    except Exception as exc:
        logger.warning("OpenRouter error: %s", exc)
    return None


async def query_ai(sid: str, prompt: str, system: Optional[str] = None) -> str:
    user_stats[sid] = user_stats.get(sid, 0) + 1
    sys_msg = system if system is not None else SYSTEM_PROMPT
    preferred = user_model.get(sid, "groq")

    if preferred == "groq":
        order = [("Groq", query_groq), ("Gemini", query_gemini), ("OpenRouter", query_openrouter)]
    elif preferred == "gemini":
        order = [("Gemini", query_gemini), ("Groq", query_groq), ("OpenRouter", query_openrouter)]
    else:
        order = [("OpenRouter", query_openrouter), ("Groq", query_groq), ("Gemini", query_gemini)]

    for name, engine in order:
        result = await engine(sid, prompt, sys_msg)
        if result:
            return result
        logger.info("%s sin resultado, probando siguiente motor...", name)

    ninguno = groq_client is None and gemini_client is None and not OPENROUTER_API_KEY
    if ninguno:
        return "⚠️ Sin motores de IA activos. Configura al menos una clave API en el panel de configuración."
    return "⚠️ Los motores están ocupados en este momento. Intenta de nuevo en unos segundos."


# ══════════════════════════════════════════════════════════════════════════════
#  GENERADOR DE IMAGENES (Pollinations AI)
# ══════════════════════════════════════════════════════════════════════════════
async def generate_image(prompt: str) -> Optional[bytes]:
    enhanced = (
        "academic educational scientific illustration, university textbook style, "
        "professional labeled diagram, clean white background, high detail, "
        "no human faces, spanish labels: " + prompt
    )
    encoded = urllib.parse.quote(enhanced)
    url = (
        f"https://image.pollinations.ai/prompt/{encoded}"
        "?width=1024&height=1024&nologo=true&enhance=true&model=flux"
    )
    try:
        async with httpx.AsyncClient(timeout=90.0, follow_redirects=True) as client:
            resp = await client.get(url)
            ct = resp.headers.get("content-type", "")
            if resp.status_code == 200 and "image" in ct:
                return resp.content
    except Exception as exc:
        logger.error("Pollinations error: %s", exc)
    return None


# ══════════════════════════════════════════════════════════════════════════════
#  NOTICIAS RSS
# ══════════════════════════════════════════════════════════════════════════════
_NOTICIAS_FEEDS: Dict[str, Dict[str, str]] = {
    "peru": {
        "El Comercio": "https://rss.elcomercio.pe/rss/portada",
        "RPP Noticias": "https://rss.rpp.pe/peru/actualidad/feed",
        "La Republica": "https://larepublica.pe/feed",
        "Andina": "https://andina.pe/agencia/rss.aspx",
        "Gestion": "https://gestion.pe/feed",
    },
    "internacional": {
        "BBC Mundo": "https://feeds.bbci.co.uk/mundo/rss.xml",
        "CNN en Espanol": "https://cnnespanol.cnn.com/feed/",
        "El Pais": "https://feeds.elpais.com/mrss-s/pages/ep/site/elpais.com/portada",
        "France 24": "https://www.france24.com/es/rss",
        "DW Espanol": "https://rss.dw.com/rdf/rss-spa-all",
    },
    "ciencia": {
        "Science Daily": "https://www.sciencedaily.com/rss/all.xml",
        "Nature": "https://www.nature.com/nature.rss",
        "El Pais Ciencia": "https://feeds.elpais.com/mrss-s/pages/ep/site/elpais.com/section/ciencia/portada",
    },
}


async def _fetch_rss(url: str, max_items: int = 5) -> List[Dict[str, str]]:
    items: List[Dict[str, str]] = []
    try:
        async with httpx.AsyncClient(timeout=15.0, follow_redirects=True) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; NexusBot/9.0)"})
            if resp.status_code != 200:
                return items
            text = resp.text
        text_clean = re.sub(r' xmlns[^"]*"[^"]*"', "", text)
        text_clean = re.sub(r"<\?xml[^>]*\?>", "", text_clean)
        try:
            root = ET.fromstring(text_clean)
        except ET.ParseError:
            return items
        entries = root.findall(".//item")
        if not entries:
            entries = root.findall(".//{http://www.w3.org/2005/Atom}entry")
        for entry in entries[:max_items]:
            titulo = (entry.findtext("title") or entry.findtext("{http://www.w3.org/2005/Atom}title") or "")
            link = entry.findtext("link") or ""
            if not link:
                link_tag = entry.find("{http://www.w3.org/2005/Atom}link")
                link = (link_tag.get("href") or "") if link_tag is not None else ""
            fecha = (entry.findtext("pubDate") or entry.findtext("{http://www.w3.org/2005/Atom}updated") or "")
            titulo = re.sub(r"<[^>]+>", "", titulo).strip()
            link = link.strip()
            if titulo:
                items.append({"titulo": titulo, "link": link, "fecha": fecha[:16]})
    except Exception as exc:
        logger.warning("RSS error %s: %s", url, exc)
    return items


async def get_noticias(categoria: str, sid: str) -> str:
    feeds = _NOTICIAS_FEEDS.get(categoria, {})
    tasks = [_fetch_rss(url, 3) for url in feeds.values()]
    resultados = await asyncio.gather(*tasks, return_exceptions=True)

    fuentes_con_items: List[Tuple[str, List[Dict[str, str]]]] = []
    for fuente, resultado in zip(feeds.keys(), resultados):
        if isinstance(resultado, list) and resultado:
            fuentes_con_items.append((fuente, resultado))

    if not fuentes_con_items:
        prompt = (
            f"El usuario quiere noticias de: {categoria}. "
            f"Proporciona un resumen de los temas más relevantes sobre "
            f"{'el Perú' if categoria == 'peru' else categoria}, "
            "con fechas aproximadas. Máximo 10 noticias con fuente y contexto breve."
        )
        return await query_ai(sid, prompt)

    iconos = {"peru": "🇵🇪", "internacional": "🌍", "ciencia": "🔬"}
    titulos = {
        "peru": "PERÚ — NOTICIAS NACIONALES",
        "internacional": "MUNDO — NOTICIAS INTERNACIONALES",
        "ciencia": "CIENCIA Y TECNOLOGÍA",
    }
    icono = iconos.get(categoria, "📰")
    titulo = titulos.get(categoria, categoria.upper())
    now_str = _now_lima().strftime("%d/%m/%Y %H:%M")

    parts: List[str] = [f"## {icono} {titulo}\n"]
    for fuente, items in fuentes_con_items:
        parts.append(f"\n### 📌 {fuente}")
        for item in items:
            t = item["titulo"][:115]
            l = item["link"]
            parts.append(f"- [{t}]({l})" if l else f"- {t}")

    parts.append(f"\n\n_🕐 Actualizado: {now_str} hora Lima_")
    return "\n".join(parts)


# ══════════════════════════════════════════════════════════════════════════════
#  PROMPTS DE COMANDOS
# ══════════════════════════════════════════════════════════════════════════════
COMMAND_PROMPTS: Dict[str, Any] = {
    "informe": lambda t: (
        f"Redacta un informe academico COMPLETO y formal sobre: {t}\n\n"
        "ESTRUCTURA (minimo 900 palabras):\n"
        "1. TITULO formal\n2. RESUMEN EJECUTIVO (150 palabras)\n"
        "3. INTRODUCCION (contexto, justificacion, objetivos)\n"
        "4. MARCO TEORICO (corrientes, autores, citas APA 7)\n"
        "5. DESARROLLO (datos empiricos, evidencia)\n"
        "6. DISCUSION (implicaciones)\n7. CONCLUSIONES (3-5 numeradas)\n"
        "8. RECOMENDACIONES (3-5)\n"
        "9. REFERENCIAS (minimo 8 fuentes APA 7, 2018-2026)"
    ),
    "ensayo": lambda t: (
        f"Ensayo academico universitario COMPLETO sobre: {t}\n\n"
        "1. TITULO formal\n2. INTRODUCCION con tesis clara\n"
        "3. DESARROLLO: 3 argumentos con evidencia y citas APA 7\n"
        "4. CONTRAPUNTO critico y refutacion\n5. POSTURA CRITICA propia\n"
        "6. CONCLUSION: sintesis y proyeccion\n"
        "7. REFERENCIAS: minimo 7 fuentes APA 7 (2018-2026)\nMinimo 750 palabras."
    ),
    "introduccion": lambda t: (
        f"Introduccion academica universitaria formal para: {t}\n\n"
        "1. GANCHO: dato estadistico o pregunta retorica\n"
        "2. CONTEXTUALIZACION: antecedentes y relevancia\n"
        "3. JUSTIFICACION\n4. PLANTEAMIENTO DEL PROBLEMA\n"
        "5. OBJETIVOS: general y especificos\n6. HIPOTESIS central\n"
        "7. ESTRUCTURA DEL TRABAJO\nMinimo 320 palabras. 3 citas APA 7."
    ),
    "conclusion": lambda t: (
        f"Conclusion academica universitaria formal para: {t}\n\n"
        "1. SINTESIS de puntos principales\n2. CONFIRMACION DE TESIS\n"
        "3. HALLAZGOS CLAVE (3-4 numerados)\n4. LIMITACIONES\n"
        "5. IMPLICACIONES practicas y teoricas\n6. RECOMENDACIONES futuras\n"
        "7. REFLEXION FINAL\nMinimo 280 palabras."
    ),
    "parafrasear": lambda t: (
        f"Parafrasea con MAXIMO rigor anti-plagio:\n\n---\n{t}\n---\n\n"
        "1. TEXTO PARAFRASEADO\n2. ESTRATEGIAS ANTI-PLAGIO APLICADAS\n"
        "3. PORCENTAJE ESTIMADO DE ORIGINALIDAD"
    ),
    "carta": lambda t: (
        f"Redacta una carta profesional para: {t}\n\n"
        "Incluye: Lugar y fecha (Lima, Peru), destinatario, saludo, "
        "cuerpo completo, despedida formal, espacio para firma. "
        "Tono profesional. En espanol peruano."
    ),
    "abstract": lambda t: (
        f"Abstract academico APA 7 sobre: {t}\n\n"
        "Un parrafo (150-250 palabras): objetivo, metodologia, resultados, conclusiones.\n"
        "PALABRAS CLAVE: 5 en espanol\nKeywords: 5 en ingles\n"
        "VERSION EN INGLES traducida formalmente."
    ),
    "tesis": lambda t: (
        f"Planteamiento COMPLETO de tesis universitaria sobre: {t}\n\n"
        "1. TITULO TENTATIVO\n2. PROBLEMA DE INVESTIGACION\n"
        "3. JUSTIFICACION (teorica, practica, metodologica)\n"
        "4. OBJETIVOS (Taxonomia de Bloom)\n5. HIPOTESIS\n"
        "6. VARIABLES con indicadores\n7. MARCO TEORICO resumido\n"
        "8. METODOLOGIA sugerida\n9. REFERENCIAS BASE (8 APA 7)"
    ),
    "resumen": lambda t: (
        f"Resume academicamente:\n\n---\n{t}\n---\n\n"
        "1. IDEAS NUCLEARES (bullet jerarquicos)\n"
        "2. ARGUMENTO CENTRAL (2-3 oraciones)\n"
        "3. DATOS Y EVIDENCIAS clave\n4. PALABRAS CLAVE\n"
        "5. PREGUNTA REFLEXIVA para profundizar\nResumen = 30% del original."
    ),
    "mapa": lambda t: (
        f"Mapa conceptual COMPLETO sobre: {t}\n\n"
        f"🔷 [{t.upper()}]\n"
        "├── 📌 CATEGORIA 1\n│   ├── 🔹 Subcategoria 1.1\n│   └── 🔹 Subcategoria 1.2\n"
        "├── 📌 CATEGORIA 2\n│   └── ...\nMinimo 6 categorias.\n"
        "🔗 RELACIONES CONCEPTUALES\n👤 AUTORES CLAVE"
    ),
    "comparar": lambda t: (
        f"Comparacion academica exhaustiva: {t}\n\n"
        "1. TABLA COMPARATIVA (10+ criterios)\n2. ANALISIS PROFUNDO\n"
        "3. SIMILITUDES FUNDAMENTALES\n4. DIFERENCIAS CRITICAS\n"
        "5. CONTEXTO HISTORICO\n6. VIGENCIA ACTUAL\n7. CUANDO USAR CADA UNO\n"
        "8. CONCLUSION CRITICA\n9. REFERENCIAS APA 7 (5+)"
    ),
    "esquema": lambda t: (
        f"Esquema de estudio COMPLETO sobre: {t}\n\n"
        "I. TEMA PRINCIPAL\n   A. Subtema\n      1. Concepto: definicion exacta\n"
        "Minimo 6 temas principales.\n\n"
        "⭐ 7 DATOS CRITICOS para el examen\n"
        "👤 AUTORES INDISPENSABLES\n📅 FECHAS CLAVE\n❓ 5 PREGUNTAS PROBABLES"
    ),
    "critica": lambda t: (
        f"Analisis critico academico COMPLETO de: {t}\n\n"
        "1. DESCRIPCION y contexto\n2. APORTES POSITIVOS\n"
        "3. LIMITACIONES Y DEBILIDADES\n4. CRITICAS DE OTROS AUTORES\n"
        "5. VIGENCIA ACTUAL\n6. POSTURA CRITICA PROPIA\n"
        "7. INFLUENCIA en la disciplina\n8. REFERENCIAS APA 7 (6+)"
    ),
    "argumentos": lambda t: (
        f"Argumentos academicos sobre: {t}\n\n"
        "A. A FAVOR (5+): argumento + evidencia + cita APA 7\n\n"
        "B. EN CONTRA (4+): contraargumento + evidencia\n\n"
        "C. POSTURA SINTETICA integradora\n\n"
        "D. FALACIAS LOGICAS a evitar\n\n"
        "E. COMO DEFENDER LA POSTURA en debate\n\nREFERENCIAS APA 7 (7+)"
    ),
    "buscar": lambda t: (
        f"Revision bibliografica academica sobre: {t}\n\n"
        "1. ESTUDIOS EMPIRICOS RECIENTES (2022-2026)\n"
        "2. META-ANALISIS Y REVISIONES SISTEMATICAS\n"
        "3. LIBROS Y MANUALES CLAVE\n4. AUTORES REFERENTES\n"
        "5. TENDENCIAS ACTUALES\n6. PALABRAS CLAVE para bases de datos\n"
        "7. BASES DE DATOS con URLs: Scopus, WoS, PubMed, SciELO, Redalyc\n"
        "8. REVISTAS ESPECIALIZADAS (con ISSN)\n"
        "Solo cita fuentes con certeza. Formato APA 7."
    ),
    "definir": lambda t: (
        f"Define academicamente: {t}\n\n"
        "1. DEFINICION OFICIAL (DSM-5, RAE, APA segun corresponda)\n"
        "2. ETIMOLOGIA\n3. DEFINICIONES POR CORRIENTE TEORICA\n"
        "4. CARACTERISTICAS Y CRITERIOS\n5. DISTINCION de conceptos similares\n"
        "6. EJEMPLO PRACTICO\n7. EVOLUCION HISTORICA\n8. REFERENCIAS APA 7"
    ),
    "autor": lambda t: (
        f"Informacion academica completa sobre: {t}\n\n"
        "1. DATOS BIOGRAFICOS\n2. CORRIENTE TEORICA\n"
        "3. OBRAS PRINCIPALES\n4. TEORIA O MODELO MAS CONOCIDO\n"
        "5. APORTES A LA CIENCIA\n6. CRITICAS Y LIMITACIONES\n"
        "7. COMO CITARLO EN APA 7\n8. AUTORES RELACIONADOS"
    ),
    "revista": lambda t: (
        f"Revistas cientificas del area: {t}\n\n"
        "Para cada revista: nombre, ISSN, factor de impacto, indexacion, URL.\n\n"
        "A) INTERNACIONALES ALTO IMPACTO Q1-Q2 (6+)\n"
        "B) LATINOAMERICANAS INDEXADAS (4+)\n"
        "C) ACCESO ABIERTO GRATUITO (3+)\n"
        "D) BASES DE DATOS para el area\nE) TIPS PARA BUSCAR Y PUBLICAR"
    ),
    "hipotesis": lambda t: (
        f"Formula hipotesis de investigacion para: {t}\n\n"
        "1. HIPOTESIS DE INVESTIGACION (Hi)\n2. HIPOTESIS NULA (Ho)\n"
        "3. HIPOTESIS ALTERNATIVA (Ha)\n4. VARIABLES con definiciones\n"
        "5. TIPO DE INVESTIGACION sugerida\n6. DISENO\n"
        "7. INSTRUMENTOS sugeridos\n8. JUSTIFICACION con 3 referencias APA 7"
    ),
    "metodologia": lambda t: (
        f"Seccion de METODOLOGIA completa para: {t}\n\n"
        "1. TIPO DE INVESTIGACION\n2. NIVEL\n3. DISENO\n"
        "4. POBLACION Y MUESTRA\n5. TECNICAS E INSTRUMENTOS\n"
        "6. VALIDEZ Y CONFIABILIDAD\n7. PROCEDIMIENTO\n"
        "8. ANALISIS ESTADISTICO\n9. CONSIDERACIONES ETICAS\n"
        "10. REFERENCIAS APA 7 (5+)"
    ),
    "estadistica": lambda t: (
        f"Asistencia estadistica para: {t}\n\n"
        "1. TIPO DE ANALISIS recomendado\n2. ESTADISTICOS DESCRIPTIVOS\n"
        "3. PRUEBA ESTADISTICA adecuada\n4. SUPUESTOS a verificar\n"
        "5. COMO INTERPRETAR resultados\n6. COMO REPORTAR en APA 7\n"
        "7. SOFTWARE recomendado (SPSS, R, JASP)\n8. ERRORES FRECUENTES"
    ),
    "apa": lambda t: (
        f"Genera cita APA 7 con los datos: {t}\n\n"
        "1. REFERENCIA COMPLETA APA 7\n2. CITA DIRECTA en el texto\n"
        "3. CITA PARAFRASEADA\n4. CITA CON 3+ AUTORES\n"
        "5. DATOS FALTANTES marcados como [DATO REQUERIDO]\n"
        "6. TIPO DE FUENTE y reglas APA 7"
    ),
    "corregirapa": lambda t: (
        f"Analiza y corrige en APA 7:\n\nCITA: {t}\n\n"
        "1. ERRORES ENCONTRADOS con explicacion\n2. REGLA APA 7 VIOLADA\n"
        "3. CITA CORREGIDA completa\n4. CITA EN TEXTO: directa y parafraseada\n"
        "5. CONSEJO para no repetir el error"
    ),
    "bibliografia": lambda t: (
        f"LISTA DE REFERENCIAS APA 7 sobre: {t}\n\n"
        "Minimo 12 referencias REALES en orden ALFABETICO.\n"
        "LIBROS Y MANUALES:\n\nARTICULOS CIENTIFICOS (2019-2026):\n\n"
        "RECURSOS WEB INSTITUCIONALES:\n\n"
        "Solo fuentes con certeza. No inventes DOIs."
    ),
    "refs": lambda t: (
        f"Lista rapida de 8 referencias APA 7 sobre: {t}\n\n"
        "Para cada referencia: APA 7 completo, DOI si disponible, aporte en 1 linea.\n"
        "Orden: mas reciente a mas antigua. SOLO fuentes indexadas."
    ),
    "preguntas": lambda t: (
        f"Banco de PREGUNTAS DE EXAMEN universitario sobre: {t}\n\n"
        "A) DESARROLLO (5): pregunta + criterios + respuesta modelo\n\n"
        "B) OPCION MULTIPLE (10): enunciado + 4 alternativas + correcta + explicacion\n\n"
        "C) VERDADERO/FALSO (5): afirmacion + V/F + justificacion\n\n"
        "D) DEFINICION (5 conceptos)\nE) CASO PRACTICO\nF) TEMAS MAS PROBABLES (top 5)"
    ),
    "flashcards": lambda t: (
        f"15 FLASHCARDS sobre: {t}\n\n"
        "🃏 TARJETA [N]\n❓ PREGUNTA: ...\n✅ RESPUESTA: ...\n"
        "💡 TRUCO MNEMONICO: ..."
    ),
    "plan": lambda t: (
        f"Plan de estudio DETALLADO para: {t}\n\n"
        "1. DIAGNOSTICO\n2. DISTRIBUCION DIARIA con temas y horas\n"
        "3. TECNICAS DE ESTUDIO\n4. MATERIAL necesario\n"
        "5. CRONOGRAMA de repasos\n6. CHECKLIST de temas\n"
        "7. ESTRATEGIA para el dia del examen"
    ),
    "tecnicas": lambda t: (
        f"Tecnicas de aprendizaje para: {t}\n\n"
        "1. LECTURA\n2. MEMORIZACION\n3. ORGANIZACION\n4. PRACTICA\n"
        "5. REPASO ESPACIADO\n6. RUTINA OPTIMA\n7. ERRORES FRECUENTES\n"
        "8. APPS recomendadas\n9. COMO MEDIR PROGRESO"
    ),
    "pomodoro": lambda t: (
        f"Guia POMODORO universitaria{f' para: {t}' if t else ''}.\n\n"
        "1. TECNICA POMODORO\n2. CONFIGURACION para universitarios\n"
        "3. PREPARAR EL ESPACIO\n4. DURANTE el pomodoro\n5. DESCANSOS\n"
        "6. REGISTRAR AVANCE\n7. VARIACIONES\n8. APPS GRATUITAS\n9. TABLA SEMANAL"
    ),
    "glosario": lambda t: (
        f"GLOSARIO ACADEMICO de 25 terminos mas importantes sobre: {t}\n\n"
        "Para cada termino:\n📝 TERMINO: definicion academica\n"
        "🔤 Sinonimos\n📚 Autor que lo definio\n\n"
        "Orden alfabetico. Al final: 5 TERMINOS CLAVE indispensables."
    ),
    "objetivos": lambda t: (
        f"Objetivos de investigacion/aprendizaje sobre: {t}\n\n"
        "TAXONOMIA DE BLOOM:\n"
        "1️⃣ RECORDAR (identificar, nombrar)\n"
        "2️⃣ COMPRENDER (explicar, describir)\n"
        "3️⃣ APLICAR (utilizar, resolver)\n"
        "4️⃣ ANALIZAR (descomponer, examinar)\n"
        "5️⃣ EVALUAR (juzgar, argumentar)\n"
        "6️⃣ CREAR (diseñar, formular)\n\n"
        "OBJETIVO GENERAL integrador\nCRITERIOS DE EVALUACION por nivel"
    ),
    "tesauro": lambda t: (
        f"Tesauro academico para: {t}\n\n"
        "1. DEFINICION ACADEMICA\n2. SINONIMOS ACADEMICOS (10+)\n"
        "3. TERMINOS RELACIONADOS\n4. ANTONIMOS\n"
        "5. HIPERONIMOS e HIPONIMOS\n6. USO EN DISTINTAS DISCIPLINAS\n"
        "7. EJEMPLOS EN REDACCION ACADEMICA\n8. PALABRAS A EVITAR\n"
        "9. EQUIVALENTES EN INGLES para bases de datos"
    ),
    "osint": lambda t: (
        f"Verificador de fuentes academicas.\n\nDATO A VERIFICAR: {t}\n\n"
        "1. VEREDICTO INICIAL: Confiable / Dudosa / No verificable\n"
        "2. ANALISIS DE CREDIBILIDAD\n3. DONDE VERIFICAR (URLs academicas)\n"
        "4. SEÑALES DE ALERTA\n5. FUENTES ALTERNATIVAS CONFIABLES\n"
        "6. RECOMENDACION FINAL"
    ),
    "traducir": lambda t: (
        f"Traduccion academica:\n\n---\n{t}\n---\n\n"
        "1. DETECCION del idioma original\n2. TRADUCCION (formal y academica)\n"
        "3. TERMINOS CLAVE con equivalente tecnico\n4. NOTAS DE TRADUCCION"
    ),
    "corregir": lambda t: (
        f"Corrige el texto academico:\n\n---\n{t}\n---\n\n"
        "1. TEXTO CORREGIDO (completo)\n"
        "2. ERRORES: ortograficos, gramaticales, puntuacion, estilo\n"
        "3. MEJORAS DE ESTILO\n4. VOCABULARIO ACADEMICO alternativo"
    ),
}


# ══════════════════════════════════════════════════════════════════════════════
#  FLASK APP
# ══════════════════════════════════════════════════════════════════════════════
from flask import Flask, request, jsonify, Response
from flask_cors import CORS

app = Flask(__name__)
CORS(app)


def run_async(coro):
    """Ejecuta una corutina desde contexto síncrono de Flask (compatible con threaded=True)."""
    try:
        loop = asyncio.get_running_loop()
        # Si ya hay un loop corriendo (no debería en Flask sync), usar thread
        import concurrent.futures
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
            future = pool.submit(asyncio.run, coro)
            return future.result()
    except RuntimeError:
        # No hay loop corriendo — caso normal en hilos de Flask
        return asyncio.run(coro)


# ── API: Estado de motores ─────────────────────────────────────────────────
@app.route("/api/status")
def api_status():
    motores = []
    if groq_client:
        motores.append({"id": "groq", "nombre": "Groq llama-3.3-70b", "icono": "⚡"})
    if gemini_client:
        motores.append({"id": "gemini", "nombre": "Gemini 2.0 Flash", "icono": "🌟"})
    if OPENROUTER_API_KEY:
        motores.append({"id": "openrouter", "nombre": "OpenRouter Mistral-7b", "icono": "🔀"})
    return jsonify({
        "version": BOT_VERSION,
        "motores": motores,
        "motores_activos": len(motores),
    })


# ── API: Configurar claves ─────────────────────────────────────────────────
@app.route("/api/config", methods=["POST"])
def api_config():
    global GROQ_API_KEY, GEMINI_API_KEY, OPENROUTER_API_KEY
    data = request.get_json(force=True)
    changed = False
    if data.get("groq_key"):
        GROQ_API_KEY = data["groq_key"].strip()
        changed = True
    if data.get("gemini_key"):
        GEMINI_API_KEY = data["gemini_key"].strip()
        changed = True
    if data.get("openrouter_key"):
        OPENROUTER_API_KEY = data["openrouter_key"].strip()
        changed = True
    if changed:
        _init_clients()
    return jsonify({"ok": True, "mensaje": "Configuración actualizada correctamente."})


# ── API: Chat ─────────────────────────────────────────────────────────────
@app.route("/api/chat", methods=["POST"])
def api_chat():
    data = request.get_json(force=True)
    sid = data.get("session_id", "default")
    mensaje = data.get("mensaje", "").strip()
    modo = data.get("modo", user_mode.get(sid, "libre"))
    modelo = data.get("modelo", user_model.get(sid, "groq"))

    user_mode[sid] = modo
    user_model[sid] = modelo

    if not mensaje:
        return jsonify({"error": "Mensaje vacío"}), 400

    sys_msg = MODE_SYSTEMS.get(modo, MODE_SYSTEMS["libre"])

    # Modo clínico
    if modo == "clinico" and sid in paciente_sessions:
        ficha = paciente_sessions[sid].get("ficha", "")
        sys_msg += f"\n\nFICHA CLINICA (no revelar):\n{ficha}"

    # Inyectar contexto de documento si existe
    ctx_doc = doc_context.get(sid)
    if ctx_doc:
        hist_refs = conversation_history.get(sid, [])
        tiene_doc = any("[Documento:" in m.get("content", "") for m in hist_refs[-6:])
        if tiene_doc and ctx_doc.get("texto"):
            sys_msg += (
                f"\n\n[CONTEXTO: El usuario cargó '{ctx_doc['filename']}' "
                f"({ctx_doc['tipo']}). Contenido:\n{ctx_doc['texto'][:6000]}\n"
                "Usa este contenido para responder.]"
            )

    respuesta = run_async(query_ai(sid, mensaje, sys_msg))

    if modo == "clinico":
        respuesta = f"🛋️ **[Paciente]:** _{respuesta}_"

    return jsonify({"respuesta": respuesta, "stats": user_stats.get(sid, 0)})


# ── API: Comando ──────────────────────────────────────────────────────────
@app.route("/api/comando", methods=["POST"])
def api_comando():
    data = request.get_json(force=True)
    sid = data.get("session_id", "default")
    cmd = data.get("comando", "").strip().lower()
    args = data.get("args", "").strip()
    modelo = data.get("modelo", user_model.get(sid, "groq"))
    user_model[sid] = modelo

    # Comandos especiales
    if cmd == "nuevo":
        conversation_history.pop(sid, None)
        paciente_sessions.pop(sid, None)
        doc_context.pop(sid, None)
        if user_mode.get(sid) == "clinico":
            user_mode[sid] = "libre"
        return jsonify({"respuesta": "🗑️ **Historial borrado.** Nueva conversación iniciada.", "stats": 0})

    if cmd == "hora":
        now = _now_lima()
        dias = ["Lunes","Martes","Miercoles","Jueves","Viernes","Sabado","Domingo"]
        meses = ["enero","febrero","marzo","abril","mayo","junio","julio","agosto","septiembre","octubre","noviembre","diciembre"]
        return jsonify({"respuesta": (
            f"## 🕐 Fecha y hora actual\n\n"
            f"📅 **{dias[now.weekday()]}, {now.day} de {meses[now.month-1]} de {now.year}**\n\n"
            f"⏰ **{now.strftime('%H:%M:%S')}**\n\n"
            f"🌍 Zona: America/Lima (UTC-5)"
        )})

    if cmd == "stats":
        tareas = user_tasks.get(sid, [])
        hechas = sum(1 for t in tareas if t.get("done"))
        return jsonify({"respuesta": (
            f"## 📊 Tus estadísticas\n\n"
            f"💬 Mensajes enviados: **{user_stats.get(sid, 0)}**\n\n"
            f"⚙️ Modo: **{user_mode.get(sid, 'libre')}**\n\n"
            f"🤖 Motor: **{user_model.get(sid, 'groq').upper()}**\n\n"
            f"✅ Tareas completadas: **{hechas}/{len(tareas)}**\n\n"
            f"🎯 Meta: **{user_goals.get(sid, 'Sin definir')}**"
        )})

    if cmd == "tarea":
        if not args:
            return jsonify({"respuesta": "✏️ **Uso:** Escribe la descripción de la tarea en el campo de argumentos."})
        tarea = {"id": len(user_tasks[sid]) + 1, "texto": args, "done": False, "fecha": _now_lima().strftime("%d/%m %H:%M")}
        user_tasks[sid].append(tarea)
        return jsonify({"respuesta": f"✅ **Tarea #{tarea['id']} agregada**\n\n📝 {args}"})

    if cmd == "tareas":
        tasks = user_tasks.get(sid, [])
        if not tasks:
            return jsonify({"respuesta": "📋 No tienes tareas. Usa el comando **tarea** para agregar."})
        pendientes = [t for t in tasks if not t["done"]]
        completadas = [t for t in tasks if t["done"]]
        msg = "## 📋 Tus Tareas\n\n"
        if pendientes:
            msg += "### ⏳ Pendientes\n"
            for t in pendientes:
                msg += f"- **#{t['id']}** — {t['texto']} _{t['fecha']}_\n"
        if completadas:
            msg += "\n### ✅ Completadas\n"
            for t in completadas:
                msg += f"- ~~#{t['id']} {t['texto']}~~\n"
        msg += f"\n_Total: {len(tasks)} | Pendientes: {len(pendientes)} | Completadas: {len(completadas)}_"
        return jsonify({"respuesta": msg})

    if cmd == "hecha":
        try:
            num = int(args)
            for t in user_tasks.get(sid, []):
                if t["id"] == num:
                    t["done"] = True
                    return jsonify({"respuesta": f"🎉 **Tarea #{num} completada!**\n✅ {t['texto']}"})
            return jsonify({"respuesta": f"❌ No encontré la tarea #{num}."})
        except ValueError:
            return jsonify({"respuesta": "❌ Escribe el número de tarea. Ej: `hecha` con args `2`"})

    if cmd == "objetivo":
        if not args:
            meta = user_goals.get(sid, "")
            if meta:
                return jsonify({"respuesta": f"🎯 **Tu meta actual:**\n\n_{meta}_"})
            return jsonify({"respuesta": "🎯 Escribe tu meta académica en el campo de argumentos."})
        user_goals[sid] = args
        prompt = (
            f"Ayuda a este estudiante a alcanzar: {args}\n\n"
            "1. ANALISIS DE LA META\n2. PLAN DE ACCION en 3 etapas\n"
            "3. HABITOS DIARIOS\n4. OBSTACULOS y como superarlos\n"
            "5. INDICADORES DE PROGRESO\n6. FRASE MOTIVACIONAL\n7. PRIMER PASO para HOY"
        )
        respuesta = run_async(query_ai(sid, prompt))
        return jsonify({"respuesta": f"🎯 **Meta:** _{args}_\n\n{respuesta}"})

    if cmd == "noticias":
        categoria = args.lower() if args else "peru"
        if categoria in ("mundo", "internacional", "world"):
            categoria = "internacional"
        elif categoria in ("ciencia", "tech", "tecnologia"):
            categoria = "ciencia"
        else:
            categoria = "peru"
        respuesta = run_async(get_noticias(categoria, sid))
        return jsonify({"respuesta": respuesta})

    if cmd == "paciente":
        user_mode[sid] = "clinico"
        conversation_history[sid] = []
        trastorno = args if args else ""
        prompt_caso = (
            f"Crea un caso clínico para práctica de entrevista psicológica universitaria.\n\n"
            f"{'Trastorno: ' + trastorno if trastorno else 'Elige un trastorno DSM-5 aleatorio.'}\n\n"
            "ESTRUCTURA:\n1. NOMBRE FICTICIO y datos sociodemograficos\n"
            "2. MOTIVO DE CONSULTA\n3. DIAGNOSTICO PRESUNTIVO (DSM-5)\n"
            "4. SINTOMAS PRINCIPALES\n5. HISTORIA CLINICA BREVE\n"
            "6. ESTILO DE PERSONALIDAD\n7. COMO RESPONDER en la entrevista\n"
            "8. OBJETIVOS DE EVALUACION\n\n"
            "Al final escribe exactamente:\n===INICIO_SESION===\n"
            "[primer mensaje del paciente al entrar a consulta, primera persona]"
        )
        respuesta_caso = run_async(query_ai(sid, prompt_caso, SYSTEM_PROMPT))
        if "===INICIO_SESION===" in respuesta_caso:
            partes = respuesta_caso.split("===INICIO_SESION===", 1)
            ficha = partes[0].strip()
            inicio = partes[1].strip()
        else:
            ficha = respuesta_caso
            inicio = "Hola... Me dijeron que viniera aquí. No sé muy bien para qué."

        paciente_sessions[sid] = {"ficha": ficha, "inicio": inicio, "activo": True}
        conversation_history[sid] = [{"role": "assistant", "content": inicio}]

        return jsonify({"respuesta": (
            f"## 🧠 Ficha Clínica — Solo para el terapeuta\n\n{ficha}\n\n"
            f"---\n💡 **Instrucciones:** Realiza la entrevista clínica como si fuera real. "
            f"Usa /feedback para evaluación parcial. /salirclinico para terminar.\n\n"
            f"---\n### 🛋️ [PACIENTE ENTRA A CONSULTA]\n\n_{inicio}_"
        )})

    if cmd == "salirclinico":
        sesion = paciente_sessions.get(sid)
        if not sesion:
            return jsonify({"respuesta": "No hay una sesión clínica activa."})
        historial = conversation_history.get(sid, [])
        user_mode[sid] = "libre"
        paciente_sessions.pop(sid, None)
        conversation_history.pop(sid, None)
        if len(historial) < 2:
            return jsonify({"respuesta": "✅ Sesión clínica finalizada."})
        ficha = sesion.get("ficha", "")
        intercambio = "\n".join(
            f"{'Terapeuta' if m['role']=='user' else 'Paciente'}: {m['content']}"
            for m in historial[:20]
        )
        prompt_eval = (
            f"Evalúa esta práctica de entrevista psicológica.\n\nFICHA:\n{ficha}\n\n"
            f"TRANSCRIPCION:\n{intercambio}\n\n"
            "EVALUACION:\n1. RAPPORT Y ALIANZA (0-10)\n2. TECNICAS usadas\n"
            "3. EXPLORACION DE SINTOMAS\n4. HIPOTESIS DIAGNOSTICA\n"
            "5. ERRORES Y CORRECCIONES\n6. FORTALEZAS\n"
            "7. PUNTUACION FINAL (0-20)\n8. RECOMENDACIONES"
        )
        eval_resp = run_async(query_ai(sid, prompt_eval, SYSTEM_PROMPT))
        return jsonify({"respuesta": f"## 📋 Evaluación de tu Entrevista Clínica\n\n{eval_resp}"})

    if cmd == "feedback":
        if sid not in paciente_sessions:
            return jsonify({"respuesta": "No hay sesión clínica activa. Usa el comando **paciente**."})
        historial = conversation_history.get(sid, [])
        intercambio = "\n".join(
            f"{'Terapeuta' if m['role']=='user' else 'Paciente'}: {m['content']}"
            for m in historial[-10:]
        )
        prompt = (
            f"Retroalimentación rápida sobre los últimos intercambios:\n\n{intercambio}\n\n"
            "Evalúa técnicas, aciertos y errores en máximo 200 palabras. "
            "1 sugerencia concreta para las próximas preguntas."
        )
        respuesta = run_async(query_ai(sid, prompt, SYSTEM_PROMPT))
        return jsonify({"respuesta": f"💡 **Feedback parcial:**\n\n{respuesta}"})

    if cmd == "leer":
        return jsonify({"respuesta": (
            "## 📄 Lectura de Documentos\n\n"
            "Usa el botón **📎 Adjuntar archivo** para subir tu documento.\n\n"
            "**Formatos soportados:**\n"
            "- PDF — artículos, libros, tesis\n"
            "- DOCX — documentos Word\n"
            "- XLSX — hojas de cálculo Excel\n"
            "- PPTX — presentaciones PowerPoint\n"
            "- TXT / CSV / MD — texto plano\n\n"
            "Después de subir el archivo, puedes hacerme preguntas sobre él."
        )})

    # Comandos con prompts predefinidos
    if cmd in COMMAND_PROMPTS:
        if not args and cmd not in ("pomodoro",):
            return jsonify({"respuesta": f"⚠️ Escribe el tema o texto en el campo **Argumentos** antes de ejecutar **/{cmd}**."})
        prompt_fn = COMMAND_PROMPTS[cmd]
        prompt = prompt_fn(args)
        respuesta = run_async(query_ai(sid, prompt))
        return jsonify({"respuesta": respuesta})

    return jsonify({"respuesta": f"❓ Comando `/{cmd}` no reconocido. Escribe **/ayuda** para ver todos los comandos."})


# ── API: Subir documento ──────────────────────────────────────────────────
@app.route("/api/documento", methods=["POST"])
def api_documento():
    sid = request.form.get("session_id", "default")
    pregunta = request.form.get("pregunta", "").strip()
    modelo = request.form.get("modelo", user_model.get(sid, "groq"))
    user_model[sid] = modelo

    if "archivo" not in request.files:
        return jsonify({"error": "No se recibió ningún archivo"}), 400

    archivo = request.files["archivo"]
    filename = archivo.filename or "archivo"
    mime = archivo.content_type or ""
    file_bytes = archivo.read()

    texto, tipo = extract_text_from_file(file_bytes, mime, filename)

    if not texto:
        return jsonify({"respuesta": (
            f"⚠️ No pude extraer texto de **{filename}**.\n\n"
            f"Tipo detectado: {tipo}\n"
            "Asegúrate de que el PDF no sea una imagen escaneada sin OCR."
        )})

    pregunta_final = pregunta if pregunta else (
        "Resume este documento de forma académica: identifica el tipo de texto, "
        "las ideas principales, argumentos centrales, metodología si la hay, "
        "conclusiones y palabras clave. Organiza con subtítulos claros."
    )

    prompt = (
        f"El usuario envió un archivo **{tipo}** llamado '{filename}'.\n\n"
        f"CONTENIDO:\n{'─'*40}\n{texto}\n{'─'*40}\n\n"
        f"INSTRUCCION: {pregunta_final}\n\n"
        "Responde en español con lenguaje académico."
    )

    doc_context[sid] = {"filename": filename, "tipo": tipo, "texto": texto}
    conversation_history.setdefault(sid, [])
    conversation_history[sid].append({
        "role": "user",
        "content": f"[Documento: {filename} ({tipo})] Pregunta: {pregunta_final}",
    })

    respuesta = run_async(query_ai(sid, prompt))
    conversation_history[sid].append({"role": "assistant", "content": respuesta})

    return jsonify({
        "respuesta": f"## 📄 Análisis de: _{filename}_\n\n{respuesta}\n\n_💡 Ahora puedes hacerme preguntas sobre este documento._"
    })


# ── API: Generar imagen ───────────────────────────────────────────────────
@app.route("/api/imagen", methods=["POST"])
def api_imagen():
    data = request.get_json(force=True)
    prompt_desc = data.get("prompt", "").strip()
    if not prompt_desc:
        return jsonify({"error": "Prompt vacío"}), 400

    image_bytes = run_async(generate_image(prompt_desc))
    if image_bytes:
        b64 = base64.b64encode(image_bytes).decode("utf-8")
        return jsonify({"imagen": f"data:image/jpeg;base64,{b64}", "prompt": prompt_desc})
    return jsonify({"error": "No se pudo generar la imagen"}), 500


# ── API: Cita APA desde URL ───────────────────────────────────────────────
@app.route("/api/apa_web", methods=["POST"])
def api_apa_web():
    data = request.get_json(force=True)
    sid = data.get("session_id", "default")
    url = data.get("url", "").strip()
    if not url or not url.startswith("http"):
        return jsonify({"error": "URL inválida"}), 400

    async def _fetch_and_cite():
        metadata_str = f"URL: {url}"
        try:
            async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
                resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; NexusBot/9.0)"})
                if resp.status_code == 200:
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(resp.text[:50000], "lxml")

                    def meta(name: str, prop: str = "") -> str:
                        tag = soup.find("meta", attrs={"name": name}) or soup.find("meta", attrs={"property": prop or name})
                        if not tag:
                            return ""
                        val = tag.get("content", "")
                        return str(val).strip() if val else ""

                    metadata = {
                        "titulo": meta("citation_title") or meta("", "og:title") or (str(soup.title.string).strip() if soup.title and soup.title.string else ""),
                        "autores": meta("citation_author") or meta("author"),
                        "año": meta("citation_publication_date") or meta("citation_year"),
                        "revista": meta("citation_journal_title") or meta("", "og:site_name"),
                        "doi": meta("citation_doi"),
                        "url": url,
                    }
                    metadata_str = f"URL: {url}\nMetadatos: {json.dumps(metadata, ensure_ascii=False)}"
        except Exception as exc:
            logger.warning("Error scraping %s: %s", url, exc)

        prompt = (
            f"Genera una cita APA 7 para este recurso web:\n\n{metadata_str}\n\n"
            "1. REFERENCIA COMPLETA APA 7\n2. CITA EN TEXTO: directa y parafraseada\n"
            "3. TIPO DE FUENTE\n4. DATOS FALTANTES como [DATO REQUERIDO]\n"
            "5. ADVERTENCIA si la fuente no es académica"
        )
        return await query_ai(sid, prompt)

    respuesta = run_async(_fetch_and_cite())
    return jsonify({"respuesta": f"## 🔗 Cita APA 7 desde URL\n\n{respuesta}"})


# ── Servir interfaz HTML ──────────────────────────────────────────────────
@app.route("/")
def index():
    return HTML_PAGE


# ══════════════════════════════════════════════════════════════════════════════
#  INTERFAZ HTML COMPLETA
# ══════════════════════════════════════════════════════════════════════════════
HTML_PAGE = r"""<!DOCTYPE html>
<html lang="es">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>NEXUS PRIME — Asistente Académico</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Syne:wght@400;500;600;700;800&family=JetBrains+Mono:wght@400;500&family=Inter:wght@300;400;500&display=swap" rel="stylesheet">
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/styles/github-dark.min.css">
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/marked/9.1.6/marked.min.js"></script>
<style>
:root {
  --bg: #07080f;
  --bg2: #0d0f1c;
  --bg3: #131629;
  --surface: #181b30;
  --surface2: #1e2240;
  --border: #2a2f52;
  --border2: #353b6b;
  --accent: #6c63ff;
  --accent2: #9d97ff;
  --accent3: #c4c1ff;
  --green: #00e5a0;
  --gold: #ffd166;
  --red: #ff6b6b;
  --text: #e8eaf6;
  --text2: #9fa8da;
  --text3: #5c6494;
  --radius: 16px;
  --shadow: 0 8px 32px rgba(0,0,0,0.5);
  --glow: 0 0 40px rgba(108,99,255,0.15);
}

* { margin: 0; padding: 0; box-sizing: border-box; }

body {
  background: var(--bg);
  color: var(--text);
  font-family: 'Inter', sans-serif;
  height: 100dvh;
  display: flex;
  overflow: hidden;
  font-size: 14px;
}

/* ── SIDEBAR ── */
#sidebar {
  width: 300px;
  min-width: 300px;
  background: var(--bg2);
  border-right: 1px solid var(--border);
  display: flex;
  flex-direction: column;
  overflow: hidden;
  transition: transform 0.3s ease;
}

.sidebar-header {
  padding: 24px 20px 16px;
  border-bottom: 1px solid var(--border);
  background: linear-gradient(135deg, var(--bg2) 0%, rgba(108,99,255,0.08) 100%);
}

.logo-row {
  display: flex;
  align-items: center;
  gap: 12px;
  margin-bottom: 12px;
}

.logo-icon {
  width: 42px;
  height: 42px;
  background: linear-gradient(135deg, var(--accent), #9d4edd);
  border-radius: 12px;
  display: flex;
  align-items: center;
  justify-content: center;
  font-size: 20px;
  box-shadow: 0 4px 16px rgba(108,99,255,0.4);
  flex-shrink: 0;
}

.logo-text h1 {
  font-family: 'Syne', sans-serif;
  font-size: 18px;
  font-weight: 800;
  color: var(--text);
  letter-spacing: 1px;
}

.logo-text p {
  font-size: 10px;
  color: var(--text3);
  letter-spacing: 0.5px;
  text-transform: uppercase;
}

.engine-status {
  display: flex;
  gap: 6px;
  flex-wrap: wrap;
}

.engine-tag {
  background: rgba(108,99,255,0.15);
  border: 1px solid rgba(108,99,255,0.3);
  border-radius: 20px;
  padding: 3px 10px;
  font-size: 10px;
  color: var(--accent3);
  font-family: 'JetBrains Mono', monospace;
}

.engine-tag.offline {
  background: rgba(255,107,107,0.1);
  border-color: rgba(255,107,107,0.3);
  color: #ff9999;
}

/* ── SIDEBAR TABS ── */
.sidebar-tabs {
  display: flex;
  border-bottom: 1px solid var(--border);
}

.sidebar-tab {
  flex: 1;
  padding: 10px;
  text-align: center;
  font-size: 11px;
  color: var(--text3);
  cursor: pointer;
  transition: all 0.2s;
  font-weight: 500;
  letter-spacing: 0.5px;
  text-transform: uppercase;
}

.sidebar-tab.active {
  color: var(--accent2);
  border-bottom: 2px solid var(--accent);
  background: rgba(108,99,255,0.05);
}

.sidebar-tab:hover:not(.active) { background: rgba(255,255,255,0.03); color: var(--text2); }

.tab-content { display: none; flex: 1; overflow-y: auto; }
.tab-content.active { display: block; }

/* ── COMMANDS ── */
.cmd-section { padding: 12px 16px 4px; }

.cmd-section-title {
  font-size: 10px;
  text-transform: uppercase;
  letter-spacing: 1px;
  color: var(--text3);
  font-weight: 600;
  margin-bottom: 6px;
  padding-left: 4px;
}

.cmd-btn {
  width: 100%;
  display: flex;
  align-items: center;
  gap: 10px;
  padding: 8px 12px;
  border-radius: 10px;
  border: none;
  background: transparent;
  color: var(--text2);
  cursor: pointer;
  font-size: 12.5px;
  text-align: left;
  transition: all 0.15s;
  font-family: inherit;
  margin-bottom: 2px;
}

.cmd-btn:hover {
  background: rgba(108,99,255,0.12);
  color: var(--text);
  transform: translateX(3px);
}

.cmd-btn .cmd-icon { font-size: 14px; width: 20px; text-align: center; flex-shrink: 0; }
.cmd-btn .cmd-name { font-weight: 500; }
.cmd-btn .cmd-desc { font-size: 10.5px; color: var(--text3); display: block; }

/* ── CONFIG PANEL ── */
.config-panel { padding: 16px; }

.config-title {
  font-size: 11px;
  text-transform: uppercase;
  letter-spacing: 1px;
  color: var(--text3);
  margin-bottom: 12px;
  font-weight: 600;
}

.config-group { margin-bottom: 16px; }

.config-label {
  font-size: 11px;
  color: var(--text2);
  margin-bottom: 6px;
  display: block;
  font-weight: 500;
}

.config-input {
  width: 100%;
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 8px;
  padding: 8px 12px;
  color: var(--text);
  font-size: 12px;
  font-family: 'JetBrains Mono', monospace;
  outline: none;
  transition: border 0.2s;
}

.config-input:focus { border-color: var(--accent); }
.config-input::placeholder { color: var(--text3); }

.mode-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 6px; }

.mode-btn {
  padding: 8px 6px;
  border-radius: 8px;
  border: 1px solid var(--border);
  background: var(--bg3);
  color: var(--text2);
  cursor: pointer;
  font-size: 11px;
  text-align: center;
  transition: all 0.2s;
  font-family: inherit;
  font-weight: 500;
}

.mode-btn.active {
  border-color: var(--accent);
  background: rgba(108,99,255,0.2);
  color: var(--accent3);
}

.mode-btn:hover:not(.active) { border-color: var(--border2); background: var(--surface); color: var(--text); }

.model-select {
  width: 100%;
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 8px;
  padding: 8px 12px;
  color: var(--text);
  font-size: 12px;
  outline: none;
  cursor: pointer;
  font-family: inherit;
}

.save-btn {
  width: 100%;
  padding: 10px;
  border-radius: 10px;
  border: none;
  background: linear-gradient(135deg, var(--accent), #9d4edd);
  color: white;
  font-size: 13px;
  font-weight: 600;
  cursor: pointer;
  transition: all 0.2s;
  font-family: 'Syne', sans-serif;
  letter-spacing: 0.5px;
  margin-top: 8px;
}

.save-btn:hover { opacity: 0.9; transform: translateY(-1px); box-shadow: 0 4px 16px rgba(108,99,255,0.4); }

/* ── MAIN AREA ── */
#main {
  flex: 1;
  display: flex;
  flex-direction: column;
  overflow: hidden;
  background: var(--bg);
}

/* ── TOP BAR ── */
#topbar {
  padding: 14px 20px;
  border-bottom: 1px solid var(--border);
  display: flex;
  align-items: center;
  gap: 12px;
  background: var(--bg2);
  flex-shrink: 0;
}

#menu-toggle {
  display: none;
  background: none;
  border: none;
  color: var(--text2);
  cursor: pointer;
  font-size: 20px;
  padding: 4px;
}

.topbar-title {
  font-family: 'Syne', sans-serif;
  font-weight: 700;
  font-size: 15px;
  color: var(--text);
  flex: 1;
}

.topbar-meta {
  font-size: 11px;
  color: var(--text3);
  font-family: 'JetBrains Mono', monospace;
}

.clear-btn {
  padding: 6px 14px;
  border-radius: 20px;
  border: 1px solid var(--border2);
  background: transparent;
  color: var(--text2);
  font-size: 11px;
  cursor: pointer;
  transition: all 0.2s;
  font-family: inherit;
  font-weight: 500;
}

.clear-btn:hover { background: rgba(255,107,107,0.15); border-color: var(--red); color: var(--red); }

/* ── MESSAGES ── */
#messages {
  flex: 1;
  overflow-y: auto;
  padding: 20px;
  display: flex;
  flex-direction: column;
  gap: 16px;
  scroll-behavior: smooth;
}

#messages::-webkit-scrollbar { width: 6px; }
#messages::-webkit-scrollbar-track { background: transparent; }
#messages::-webkit-scrollbar-thumb { background: var(--border); border-radius: 3px; }

/* ── WELCOME ── */
#welcome {
  display: flex;
  flex-direction: column;
  align-items: center;
  justify-content: center;
  flex: 1;
  text-align: center;
  padding: 40px 20px;
  gap: 20px;
}

.welcome-logo {
  width: 80px;
  height: 80px;
  background: linear-gradient(135deg, var(--accent), #9d4edd);
  border-radius: 24px;
  display: flex;
  align-items: center;
  justify-content: center;
  font-size: 36px;
  box-shadow: 0 8px 40px rgba(108,99,255,0.4);
  animation: float 3s ease-in-out infinite;
}

@keyframes float {
  0%, 100% { transform: translateY(0); }
  50% { transform: translateY(-8px); }
}

.welcome-title {
  font-family: 'Syne', sans-serif;
  font-size: 32px;
  font-weight: 800;
  background: linear-gradient(135deg, var(--text), var(--accent3));
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  background-clip: text;
}

.welcome-sub { color: var(--text3); font-size: 14px; line-height: 1.6; max-width: 400px; }

.quick-grid {
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: 10px;
  width: 100%;
  max-width: 560px;
}

.quick-card {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 12px;
  padding: 14px 12px;
  cursor: pointer;
  transition: all 0.2s;
  text-align: left;
}

.quick-card:hover {
  border-color: var(--accent);
  background: rgba(108,99,255,0.08);
  transform: translateY(-2px);
  box-shadow: 0 4px 16px rgba(108,99,255,0.15);
}

.quick-card-icon { font-size: 22px; margin-bottom: 6px; display: block; }
.quick-card-title { font-size: 12px; font-weight: 600; color: var(--text); }
.quick-card-desc { font-size: 10.5px; color: var(--text3); margin-top: 2px; }

/* ── MESSAGE BUBBLES ── */
.msg {
  display: flex;
  gap: 12px;
  animation: fadeSlide 0.25s ease;
  max-width: 100%;
}

@keyframes fadeSlide {
  from { opacity: 0; transform: translateY(10px); }
  to { opacity: 1; transform: translateY(0); }
}

.msg.user { flex-direction: row-reverse; }

.msg-avatar {
  width: 36px;
  height: 36px;
  border-radius: 10px;
  display: flex;
  align-items: center;
  justify-content: center;
  font-size: 16px;
  flex-shrink: 0;
  align-self: flex-start;
}

.msg.user .msg-avatar { background: linear-gradient(135deg, var(--accent), #9d4edd); }
.msg.bot .msg-avatar { background: var(--surface2); border: 1px solid var(--border); }

.msg-content {
  max-width: 75%;
  min-width: 0;
}

.msg-bubble {
  padding: 14px 18px;
  border-radius: 16px;
  line-height: 1.7;
  word-break: break-word;
}

.msg.user .msg-bubble {
  background: linear-gradient(135deg, var(--accent), #7c3aed);
  color: white;
  border-radius: 16px 4px 16px 16px;
  box-shadow: 0 4px 16px rgba(108,99,255,0.3);
}

.msg.bot .msg-bubble {
  background: var(--surface);
  border: 1px solid var(--border);
  color: var(--text);
  border-radius: 4px 16px 16px 16px;
}

/* Markdown inside bubbles */
.msg-bubble h1, .msg-bubble h2 {
  font-family: 'Syne', sans-serif;
  color: var(--accent3);
  margin: 16px 0 8px;
  font-size: 15px;
  border-bottom: 1px solid var(--border);
  padding-bottom: 6px;
}

.msg-bubble h1:first-child, .msg-bubble h2:first-child, .msg-bubble h3:first-child { margin-top: 0; }

.msg-bubble h3 {
  font-family: 'Syne', sans-serif;
  color: var(--text2);
  margin: 12px 0 6px;
  font-size: 13px;
}

.msg-bubble p { margin: 6px 0; }
.msg-bubble ul, .msg-bubble ol { margin: 8px 0 8px 20px; }
.msg-bubble li { margin: 4px 0; color: var(--text); }

.msg-bubble strong { color: var(--gold); font-weight: 600; }
.msg-bubble em { color: var(--accent3); font-style: italic; }

.msg-bubble a { color: var(--accent2); text-decoration: none; }
.msg-bubble a:hover { text-decoration: underline; color: var(--accent3); }

.msg-bubble code {
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 4px;
  padding: 1px 6px;
  font-family: 'JetBrains Mono', monospace;
  font-size: 12px;
  color: var(--green);
}

.msg-bubble pre {
  background: var(--bg3) !important;
  border: 1px solid var(--border);
  border-radius: 10px;
  padding: 14px 16px;
  overflow-x: auto;
  margin: 12px 0;
  font-size: 12px;
}

.msg-bubble pre code {
  background: none;
  border: none;
  padding: 0;
  color: inherit;
}

.msg-bubble blockquote {
  border-left: 3px solid var(--accent);
  padding-left: 14px;
  margin: 10px 0;
  color: var(--text2);
  font-style: italic;
}

.msg-bubble table {
  width: 100%;
  border-collapse: collapse;
  margin: 12px 0;
  font-size: 12px;
}

.msg-bubble th {
  background: var(--surface2);
  padding: 8px 12px;
  text-align: left;
  color: var(--accent3);
  border: 1px solid var(--border);
  font-weight: 600;
}

.msg-bubble td {
  padding: 7px 12px;
  border: 1px solid var(--border);
  color: var(--text2);
}

.msg-bubble tr:nth-child(even) td { background: rgba(255,255,255,0.02); }

.msg-time {
  font-size: 10px;
  color: var(--text3);
  margin-top: 4px;
  padding: 0 4px;
}

.msg.user .msg-time { text-align: right; }

/* ── TYPING ── */
.typing-dots {
  display: flex;
  gap: 5px;
  padding: 14px 18px;
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 4px 16px 16px 16px;
  width: fit-content;
}

.typing-dots span {
  width: 7px;
  height: 7px;
  border-radius: 50%;
  background: var(--accent);
  animation: typingDot 1.2s ease-in-out infinite;
}

.typing-dots span:nth-child(2) { animation-delay: 0.2s; }
.typing-dots span:nth-child(3) { animation-delay: 0.4s; }

@keyframes typingDot {
  0%, 60%, 100% { transform: translateY(0); opacity: 0.4; }
  30% { transform: translateY(-6px); opacity: 1; }
}

/* ── IMAGE MESSAGE ── */
.msg-image img {
  max-width: 100%;
  max-height: 400px;
  border-radius: 12px;
  border: 1px solid var(--border);
  display: block;
}

/* ── INPUT AREA ── */
#input-area {
  border-top: 1px solid var(--border);
  padding: 16px 20px;
  background: var(--bg2);
  flex-shrink: 0;
}

.input-extras {
  display: flex;
  gap: 8px;
  margin-bottom: 10px;
  flex-wrap: wrap;
}

.extras-label {
  font-size: 11px;
  color: var(--text3);
  align-self: center;
  font-weight: 500;
}

.extra-input {
  flex: 1;
  min-width: 140px;
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 8px;
  padding: 6px 12px;
  color: var(--text);
  font-size: 12px;
  outline: none;
  transition: border 0.2s;
  font-family: inherit;
}

.extra-input:focus { border-color: var(--accent); }
.extra-input::placeholder { color: var(--text3); }

.input-row {
  display: flex;
  gap: 10px;
  align-items: flex-end;
}

.attach-btn {
  padding: 11px 14px;
  border-radius: 10px;
  border: 1px solid var(--border);
  background: var(--surface);
  color: var(--text2);
  cursor: pointer;
  transition: all 0.2s;
  font-size: 16px;
  flex-shrink: 0;
}

.attach-btn:hover { background: var(--surface2); border-color: var(--border2); color: var(--text); }

#message-input {
  flex: 1;
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 12px;
  padding: 11px 16px;
  color: var(--text);
  font-size: 13.5px;
  resize: none;
  outline: none;
  font-family: 'Inter', sans-serif;
  line-height: 1.5;
  max-height: 130px;
  transition: border 0.2s;
  min-height: 44px;
}

#message-input:focus { border-color: var(--accent); box-shadow: 0 0 0 3px rgba(108,99,255,0.12); }
#message-input::placeholder { color: var(--text3); }

#send-btn {
  padding: 11px 20px;
  border-radius: 12px;
  border: none;
  background: linear-gradient(135deg, var(--accent), #9d4edd);
  color: white;
  font-size: 16px;
  cursor: pointer;
  transition: all 0.2s;
  flex-shrink: 0;
  box-shadow: 0 4px 16px rgba(108,99,255,0.3);
}

#send-btn:hover { transform: translateY(-1px); box-shadow: 0 6px 24px rgba(108,99,255,0.45); }
#send-btn:active { transform: translateY(0); }
#send-btn:disabled { opacity: 0.5; cursor: not-allowed; transform: none; }

.hint-text {
  font-size: 10.5px;
  color: var(--text3);
  margin-top: 8px;
  text-align: center;
}

/* ── NOTIFICATIONS ── */
#toast {
  position: fixed;
  bottom: 30px;
  right: 30px;
  background: var(--surface);
  border: 1px solid var(--border2);
  border-radius: 12px;
  padding: 12px 20px;
  font-size: 13px;
  color: var(--text);
  box-shadow: var(--shadow);
  z-index: 1000;
  transform: translateY(20px);
  opacity: 0;
  transition: all 0.3s;
  pointer-events: none;
}

#toast.show { transform: translateY(0); opacity: 1; }
#toast.success { border-color: var(--green); color: var(--green); }
#toast.error { border-color: var(--red); color: var(--red); }

/* ── MODAL (APA Web) ── */
.modal-overlay {
  position: fixed;
  inset: 0;
  background: rgba(0,0,0,0.6);
  backdrop-filter: blur(6px);
  z-index: 500;
  display: flex;
  align-items: center;
  justify-content: center;
  opacity: 0;
  pointer-events: none;
  transition: opacity 0.2s;
}

.modal-overlay.open { opacity: 1; pointer-events: all; }

.modal {
  background: var(--bg2);
  border: 1px solid var(--border);
  border-radius: 20px;
  padding: 28px;
  width: min(460px, 90vw);
  box-shadow: var(--shadow);
}

.modal h3 {
  font-family: 'Syne', sans-serif;
  font-size: 16px;
  margin-bottom: 16px;
  color: var(--accent3);
}

.modal-input {
  width: 100%;
  background: var(--bg3);
  border: 1px solid var(--border);
  border-radius: 10px;
  padding: 10px 14px;
  color: var(--text);
  font-size: 13px;
  outline: none;
  margin-bottom: 14px;
  font-family: 'JetBrains Mono', monospace;
  transition: border 0.2s;
}

.modal-input:focus { border-color: var(--accent); }

.modal-btns { display: flex; gap: 10px; }

.modal-ok {
  flex: 1;
  padding: 10px;
  border-radius: 10px;
  border: none;
  background: linear-gradient(135deg, var(--accent), #9d4edd);
  color: white;
  font-size: 13px;
  font-weight: 600;
  cursor: pointer;
  font-family: 'Syne', sans-serif;
}

.modal-cancel {
  padding: 10px 18px;
  border-radius: 10px;
  border: 1px solid var(--border);
  background: transparent;
  color: var(--text2);
  font-size: 13px;
  cursor: pointer;
  font-family: inherit;
}

/* ── RESPONSIVE ── */
@media (max-width: 768px) {
  #sidebar {
    position: fixed;
    top: 0; left: 0; bottom: 0;
    z-index: 200;
    transform: translateX(-100%);
    box-shadow: 4px 0 24px rgba(0,0,0,0.5);
  }
  #sidebar.open { transform: translateX(0); }
  #menu-toggle { display: block; }
  .quick-grid { grid-template-columns: repeat(2, 1fr); }
  .msg-content { max-width: 88%; }
}

/* ── SCROLLBAR SIDEBAR ── */
.tab-content::-webkit-scrollbar { width: 4px; }
.tab-content::-webkit-scrollbar-track { background: transparent; }
.tab-content::-webkit-scrollbar-thumb { background: var(--border); border-radius: 2px; }
</style>
</head>
<body>

<!-- APA Web Modal -->
<div class="modal-overlay" id="apa-modal">
  <div class="modal">
    <h3>🔗 Cita APA 7 desde URL</h3>
    <input class="modal-input" id="apa-url-input" type="url" placeholder="https://www.scielo.org/articulo...">
    <div class="modal-btns">
      <button class="modal-ok" onclick="submitApaWeb()">Generar cita</button>
      <button class="modal-cancel" onclick="closeModal()">Cancelar</button>
    </div>
  </div>
</div>

<!-- Sidebar -->
<aside id="sidebar">
  <div class="sidebar-header">
    <div class="logo-row">
      <div class="logo-icon">🤖</div>
      <div class="logo-text">
        <h1>NEXUS</h1>
        <p>Asistente Académico Elite</p>
      </div>
    </div>
    <div class="engine-status" id="engine-status">
      <span class="engine-tag">Cargando...</span>
    </div>
  </div>

  <div class="sidebar-tabs">
    <div class="sidebar-tab active" onclick="switchTab('comandos',this)">Comandos</div>
    <div class="sidebar-tab" onclick="switchTab('config',this)">Config</div>
  </div>

  <!-- COMANDOS TAB -->
  <div class="tab-content active" id="tab-comandos">

    <div class="cmd-section">
      <div class="cmd-section-title">📋 General</div>
      <button class="cmd-btn" onclick="runCmd('nuevo')"><span class="cmd-icon">🗑️</span><div><span class="cmd-name">Nuevo chat</span><span class="cmd-desc">Borra el historial</span></div></button>
      <button class="cmd-btn" onclick="runCmd('hora')"><span class="cmd-icon">🕐</span><div><span class="cmd-name">Hora Lima</span><span class="cmd-desc">Fecha y hora actual</span></div></button>
      <button class="cmd-btn" onclick="runCmd('stats')"><span class="cmd-icon">📊</span><div><span class="cmd-name">Estadísticas</span><span class="cmd-desc">Tu uso del bot</span></div></button>
      <button class="cmd-btn" onclick="runCmd('leer')"><span class="cmd-icon">📄</span><div><span class="cmd-name">Leer documento</span><span class="cmd-desc">PDF, DOCX, XLSX, PPTX</span></div></button>
      <button class="cmd-btn" onclick="runCmd('noticias','peru')"><span class="cmd-icon">🇵🇪</span><div><span class="cmd-name">Noticias Perú</span><span class="cmd-desc">Últimas noticias nacionales</span></div></button>
      <button class="cmd-btn" onclick="runCmd('noticias','internacional')"><span class="cmd-icon">🌍</span><div><span class="cmd-name">Noticias Mundo</span><span class="cmd-desc">Internacionales</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">✍️ Redacción</div>
      <button class="cmd-btn" onclick="focusCmd('informe')"><span class="cmd-icon">📝</span><div><span class="cmd-name">/informe</span><span class="cmd-desc">Informe académico APA</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('ensayo')"><span class="cmd-icon">📖</span><div><span class="cmd-name">/ensayo</span><span class="cmd-desc">Ensayo universitario</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('tesis')"><span class="cmd-icon">🎓</span><div><span class="cmd-name">/tesis</span><span class="cmd-desc">Planteamiento de tesis</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('introduccion')"><span class="cmd-icon">🚀</span><div><span class="cmd-name">/introduccion</span><span class="cmd-desc">Introducción formal</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('conclusion')"><span class="cmd-icon">🏁</span><div><span class="cmd-name">/conclusion</span><span class="cmd-desc">Conclusión estructurada</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('abstract')"><span class="cmd-icon">📋</span><div><span class="cmd-name">/abstract</span><span class="cmd-desc">Resumen bilingüe APA</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('carta')"><span class="cmd-icon">✉️</span><div><span class="cmd-name">/carta</span><span class="cmd-desc">Carta profesional</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('parafrasear')"><span class="cmd-icon">🔄</span><div><span class="cmd-name">/parafrasear</span><span class="cmd-desc">Anti-plagio 0%</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">🔬 Análisis</div>
      <button class="cmd-btn" onclick="focusCmd('resumen')"><span class="cmd-icon">📄</span><div><span class="cmd-name">/resumen</span><span class="cmd-desc">Resumen académico</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('mapa')"><span class="cmd-icon">🗺️</span><div><span class="cmd-name">/mapa</span><span class="cmd-desc">Mapa conceptual</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('comparar')"><span class="cmd-icon">⚖️</span><div><span class="cmd-name">/comparar</span><span class="cmd-desc">Tabla comparativa</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('esquema')"><span class="cmd-icon">📐</span><div><span class="cmd-name">/esquema</span><span class="cmd-desc">Esquema de estudio</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('critica')"><span class="cmd-icon">🔍</span><div><span class="cmd-name">/critica</span><span class="cmd-desc">Análisis crítico</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('argumentos')"><span class="cmd-icon">💡</span><div><span class="cmd-name">/argumentos</span><span class="cmd-desc">Pro y contra</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">📚 Investigación</div>
      <button class="cmd-btn" onclick="focusCmd('buscar')"><span class="cmd-icon">🔎</span><div><span class="cmd-name">/buscar</span><span class="cmd-desc">Revisión bibliográfica</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('definir')"><span class="cmd-icon">📖</span><div><span class="cmd-name">/definir</span><span class="cmd-desc">Definición académica</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('autor')"><span class="cmd-icon">👤</span><div><span class="cmd-name">/autor</span><span class="cmd-desc">Biografía y obra</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('revista')"><span class="cmd-icon">📰</span><div><span class="cmd-name">/revista</span><span class="cmd-desc">Revistas indexadas</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('hipotesis')"><span class="cmd-icon">💡</span><div><span class="cmd-name">/hipotesis</span><span class="cmd-desc">Formulación Hi/Ho</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('metodologia')"><span class="cmd-icon">🔬</span><div><span class="cmd-name">/metodologia</span><span class="cmd-desc">Diseño metodológico</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('estadistica')"><span class="cmd-icon">📊</span><div><span class="cmd-name">/estadistica</span><span class="cmd-desc">Análisis de datos</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('osint')"><span class="cmd-icon">🕵️</span><div><span class="cmd-name">/osint</span><span class="cmd-desc">Verificar fuentes</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">📚 Citas APA 7</div>
      <button class="cmd-btn" onclick="focusCmd('apa')"><span class="cmd-icon">📚</span><div><span class="cmd-name">/apa</span><span class="cmd-desc">Generar referencia</span></div></button>
      <button class="cmd-btn" onclick="openModal('apa-modal')"><span class="cmd-icon">🔗</span><div><span class="cmd-name">/apa_web</span><span class="cmd-desc">Cita desde URL</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('corregirapa')"><span class="cmd-icon">✏️</span><div><span class="cmd-name">/corregirapa</span><span class="cmd-desc">Corregir cita</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('bibliografia')"><span class="cmd-icon">📚</span><div><span class="cmd-name">/bibliografia</span><span class="cmd-desc">Lista de referencias</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('refs')"><span class="cmd-icon">📋</span><div><span class="cmd-name">/refs</span><span class="cmd-desc">Referencias rápidas</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">🧠 Estudio</div>
      <button class="cmd-btn" onclick="focusCmd('preguntas')"><span class="cmd-icon">❓</span><div><span class="cmd-name">/preguntas</span><span class="cmd-desc">Banco de examen</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('flashcards')"><span class="cmd-icon">🃏</span><div><span class="cmd-name">/flashcards</span><span class="cmd-desc">Tarjetas de estudio</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('plan')"><span class="cmd-icon">📅</span><div><span class="cmd-name">/plan</span><span class="cmd-desc">Plan de estudio</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('tecnicas')"><span class="cmd-icon">🧠</span><div><span class="cmd-name">/tecnicas</span><span class="cmd-desc">Técnicas de aprendizaje</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('pomodoro')"><span class="cmd-icon">🍅</span><div><span class="cmd-name">/pomodoro</span><span class="cmd-desc">Técnica Pomodoro</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('glosario')"><span class="cmd-icon">📖</span><div><span class="cmd-name">/glosario</span><span class="cmd-desc">Glosario de términos</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('objetivos')"><span class="cmd-icon">🎯</span><div><span class="cmd-name">/objetivos</span><span class="cmd-desc">Taxonomía de Bloom</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('tesauro')"><span class="cmd-icon">📝</span><div><span class="cmd-name">/tesauro</span><span class="cmd-desc">Sinónimos académicos</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">⚕️ Psicología Clínica</div>
      <button class="cmd-btn" onclick="focusCmd('paciente')"><span class="cmd-icon">🛋️</span><div><span class="cmd-name">/paciente</span><span class="cmd-desc">Simular caso DSM-5</span></div></button>
      <button class="cmd-btn" onclick="runCmd('feedback')"><span class="cmd-icon">💡</span><div><span class="cmd-name">/feedback</span><span class="cmd-desc">Feedback de entrevista</span></div></button>
      <button class="cmd-btn" onclick="runCmd('salirclinico')"><span class="cmd-icon">🚪</span><div><span class="cmd-name">/salirclinico</span><span class="cmd-desc">Terminar sesión</span></div></button>
    </div>

    <div class="cmd-section">
      <div class="cmd-section-title">🔤 Herramientas</div>
      <button class="cmd-btn" onclick="focusCmd('traducir')"><span class="cmd-icon">🌐</span><div><span class="cmd-name">/traducir</span><span class="cmd-desc">Traducción académica</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('corregir')"><span class="cmd-icon">✏️</span><div><span class="cmd-name">/corregir</span><span class="cmd-desc">Corrección gramatical</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('tarea')"><span class="cmd-icon">✅</span><div><span class="cmd-name">/tarea</span><span class="cmd-desc">Agregar tarea</span></div></button>
      <button class="cmd-btn" onclick="runCmd('tareas')"><span class="cmd-icon">📋</span><div><span class="cmd-name">/tareas</span><span class="cmd-desc">Ver tareas</span></div></button>
      <button class="cmd-btn" onclick="focusCmd('objetivo')"><span class="cmd-icon">🎯</span><div><span class="cmd-name">/objetivo</span><span class="cmd-desc">Meta académica</span></div></button>
      <div style="height: 20px;"></div>
    </div>
  </div>

  <!-- CONFIG TAB -->
  <div class="tab-content" id="tab-config">
    <div class="config-panel">
      <div class="config-title">⚙️ Configuración</div>

      <div class="config-group">
        <label class="config-label">🤖 Modo de respuesta</label>
        <div class="mode-grid">
          <button class="mode-btn active" onclick="setMode('libre',this)">💬 Libre</button>
          <button class="mode-btn" onclick="setMode('academico',this)">🎓 Académico</button>
          <button class="mode-btn" onclick="setMode('codigo',this)">💻 Código</button>
          <button class="mode-btn" onclick="setMode('creativo',this)">🎨 Creativo</button>
          <button class="mode-btn" onclick="setMode('tutor',this)">👨‍🏫 Tutor</button>
        </div>
      </div>

      <div class="config-group">
        <label class="config-label">⚡ Motor de IA preferido</label>
        <select class="model-select" id="model-select" onchange="setModel(this.value)">
          <option value="groq">⚡ Groq (llama-3.3-70b) — Rápido</option>
          <option value="gemini">🌟 Gemini 2.0 Flash — Versátil</option>
          <option value="openrouter">🔀 OpenRouter — Sin saturación</option>
        </select>
      </div>

      <hr style="border:none;border-top:1px solid var(--border);margin:16px 0">

      <div class="config-title">🔑 Claves API</div>
      <p style="font-size:11px;color:var(--text3);margin-bottom:12px;">Las claves se guardan en memoria mientras el servidor esté activo. Para persistencia, usa variables de entorno.</p>

      <div class="config-group">
        <label class="config-label">GROQ_API_KEY</label>
        <input class="config-input" id="groq-key" type="password" placeholder="gsk_...">
      </div>

      <div class="config-group">
        <label class="config-label">GEMINI_API_KEY</label>
        <input class="config-input" id="gemini-key" type="password" placeholder="AIza...">
      </div>

      <div class="config-group">
        <label class="config-label">OPENROUTER_API_KEY</label>
        <input class="config-input" id="or-key" type="password" placeholder="sk-or-...">
      </div>

      <button class="save-btn" onclick="saveConfig()">💾 Guardar y Reiniciar Motores</button>

      <hr style="border:none;border-top:1px solid var(--border);margin:16px 0">
      <p style="font-size:10px;color:var(--text3);text-align:center;">NEXUS PRIME v9.0 — Asistente Académico Universitario</p>
    </div>
  </div>
</aside>

<!-- Main -->
<main id="main">
  <div id="topbar">
    <button id="menu-toggle" onclick="toggleSidebar()">☰</button>
    <div class="topbar-title" id="topbar-title">💬 Chat Libre</div>
    <div class="topbar-meta" id="topbar-meta">0 mensajes</div>
    <button class="clear-btn" onclick="runCmd('nuevo')">🗑️ Limpiar</button>
  </div>

  <div id="messages">
    <div id="welcome">
      <div class="welcome-logo">🤖</div>
      <div class="welcome-title">NEXUS PRIME</div>
      <div class="welcome-sub">Tu asistente académico universitario de élite. Escríbeme directamente o usa los comandos del panel lateral.</div>
      <div class="quick-grid">
        <div class="quick-card" onclick="focusCmd('informe')">
          <span class="quick-card-icon">📝</span>
          <div class="quick-card-title">Informe APA</div>
          <div class="quick-card-desc">Completo y formal</div>
        </div>
        <div class="quick-card" onclick="focusCmd('buscar')">
          <span class="quick-card-icon">🔎</span>
          <div class="quick-card-title">Bibliografía</div>
          <div class="quick-card-desc">Fuentes indexadas</div>
        </div>
        <div class="quick-card" onclick="focusCmd('tesis')">
          <span class="quick-card-icon">🎓</span>
          <div class="quick-card-title">Tesis</div>
          <div class="quick-card-desc">Planteamiento completo</div>
        </div>
        <div class="quick-card" onclick="focusCmd('preguntas')">
          <span class="quick-card-icon">❓</span>
          <div class="quick-card-title">Banco de Examen</div>
          <div class="quick-card-desc">Preguntas y respuestas</div>
        </div>
        <div class="quick-card" onclick="focusCmd('paciente')">
          <span class="quick-card-icon">🛋️</span>
          <div class="quick-card-title">Caso Clínico</div>
          <div class="quick-card-desc">Práctica DSM-5</div>
        </div>
        <div class="quick-card" onclick="runCmd('noticias','peru')">
          <span class="quick-card-icon">📰</span>
          <div class="quick-card-title">Noticias Perú</div>
          <div class="quick-card-desc">Tiempo real</div>
        </div>
      </div>
    </div>
  </div>

  <div id="input-area">
    <div class="input-extras">
      <span class="extras-label">Cmd:</span>
      <input class="extra-input" id="cmd-input" placeholder="/informe, /tesis, /apa ..." style="max-width:160px;">
      <input class="extra-input" id="args-input" placeholder="Tema o argumento...">
    </div>
    <div class="input-row">
      <label class="attach-btn" title="Adjuntar documento">
        📎
        <input type="file" id="file-input" accept=".pdf,.docx,.xlsx,.pptx,.txt,.csv,.md" style="display:none" onchange="handleFileUpload(this)">
      </label>
      <textarea id="message-input" rows="1" placeholder="Escribe tu mensaje o pregunta académica..."></textarea>
      <button id="send-btn" onclick="sendMessage()">➤</button>
    </div>
    <div class="hint-text">Enter para enviar · Shift+Enter nueva línea · Los comandos del sidebar llenan los campos automáticamente</div>
  </div>
</main>

<div id="toast"></div>

<script>
// ── STATE ──────────────────────────────────────────────────────────────────
const state = {
  sessionId: 'session_' + Date.now() + '_' + Math.random().toString(36).slice(2),
  mode: 'libre',
  model: 'groq',
  msgCount: 0,
  loading: false,
};

// ── INIT ───────────────────────────────────────────────────────────────────
document.addEventListener('DOMContentLoaded', () => {
  loadStatus();
  setupTextarea();
  setupEnterKey();
  marked.setOptions({
    gfm: true,
    breaks: true,
    highlight: (code, lang) => {
      if (lang && hljs.getLanguage(lang)) {
        return hljs.highlight(code, { language: lang }).value;
      }
      return hljs.highlightAuto(code).value;
    }
  });
});

// ── STATUS ─────────────────────────────────────────────────────────────────
async function loadStatus() {
  try {
    const r = await fetch('/api/status');
    const d = await r.json();
    const el = document.getElementById('engine-status');
    if (d.motores.length === 0) {
      el.innerHTML = '<span class="engine-tag offline">⚠️ Sin motores — configura claves</span>';
    } else {
      el.innerHTML = d.motores.map(m => `<span class="engine-tag">${m.icono} ${m.nombre}</span>`).join('');
    }
  } catch(e) {
    document.getElementById('engine-status').innerHTML = '<span class="engine-tag offline">Sin conexión</span>';
  }
}

// ── TEXTAREA AUTO-RESIZE ───────────────────────────────────────────────────
function setupTextarea() {
  const ta = document.getElementById('message-input');
  ta.addEventListener('input', () => {
    ta.style.height = 'auto';
    ta.style.height = Math.min(ta.scrollHeight, 130) + 'px';
  });
}

function setupEnterKey() {
  document.getElementById('message-input').addEventListener('keydown', e => {
    if (e.key === 'Enter' && !e.shiftKey) {
      e.preventDefault();
      sendMessage();
    }
  });
}

// ── SEND MESSAGE ───────────────────────────────────────────────────────────
async function sendMessage() {
  if (state.loading) return;
  const input = document.getElementById('message-input');
  const cmdInput = document.getElementById('cmd-input');
  const argsInput = document.getElementById('args-input');

  const mensaje = input.value.trim();
  const cmd = cmdInput.value.trim().replace(/^\//, '').toLowerCase();
  const args = argsInput.value.trim();

  if (!mensaje && !cmd) return;

  hideWelcome();
  state.loading = true;
  document.getElementById('send-btn').disabled = true;

  if (cmd) {
    // Detectar comando /imagen o /image para generación de imágenes
    if ((cmd === 'imagen' || cmd === 'image') && args) {
      cmdInput.value = '';
      argsInput.value = '';
      state.loading = false;
      document.getElementById('send-btn').disabled = false;
      await generateImage(args);
      return;
    }
    // Comando
    const displayText = `/${cmd}${args ? ' ' + args : ''}`;
    addMessage(displayText, 'user');
    cmdInput.value = '';
    argsInput.value = '';
    const typing = addTyping();
    try {
      const r = await fetch('/api/comando', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ session_id: state.sessionId, comando: cmd, args, modelo: state.model })
      });
      const d = await r.json();
      removeTyping(typing);
      addMessage(d.respuesta, 'bot');
    } catch(e) {
      removeTyping(typing);
      addMessage('❌ Error de conexión. Verifica que el servidor esté activo.', 'bot');
    }
  } else {
    // Mensaje libre
    addMessage(mensaje, 'user');
    input.value = '';
    input.style.height = 'auto';
    const typing = addTyping();
    try {
      const r = await fetch('/api/chat', {
        method: 'POST',
        headers: {'Content-Type': 'application/json'},
        body: JSON.stringify({ session_id: state.sessionId, mensaje, modo: state.mode, modelo: state.model })
      });
      const d = await r.json();
      removeTyping(typing);
      addMessage(d.respuesta, 'bot');
      updateMeta(d.stats);
    } catch(e) {
      removeTyping(typing);
      addMessage('❌ Error de conexión. Verifica que el servidor esté activo.', 'bot');
    }
  }

  state.loading = false;
  document.getElementById('send-btn').disabled = false;
}

// ── ADD MESSAGE ────────────────────────────────────────────────────────────
function addMessage(text, role) {
  state.msgCount++;
  updateMeta(state.msgCount);

  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.className = `msg ${role}`;

  const avatar = document.createElement('div');
  avatar.className = 'msg-avatar';
  avatar.textContent = role === 'user' ? '👤' : '🤖';

  const content = document.createElement('div');
  content.className = 'msg-content';

  const bubble = document.createElement('div');
  bubble.className = 'msg-bubble';

  if (role === 'bot') {
    bubble.innerHTML = marked.parse(text || '');
    // Highlight code blocks
    bubble.querySelectorAll('pre code').forEach(el => {
      try { hljs.highlightElement(el); } catch(e) {}
    });
    // Open links in new tab
    bubble.querySelectorAll('a').forEach(a => { a.target = '_blank'; a.rel = 'noopener'; });
  } else {
    bubble.textContent = text;
  }

  const time = document.createElement('div');
  time.className = 'msg-time';
  time.textContent = new Date().toLocaleTimeString('es-PE', {hour:'2-digit', minute:'2-digit'});

  content.appendChild(bubble);
  content.appendChild(time);
  div.appendChild(avatar);
  div.appendChild(content);
  msgs.appendChild(div);
  msgs.scrollTop = msgs.scrollHeight;
  return div;
}

// ── TYPING INDICATOR ───────────────────────────────────────────────────────
function addTyping() {
  const msgs = document.getElementById('messages');
  const div = document.createElement('div');
  div.className = 'msg bot';
  div.id = 'typing-indicator';
  const avatar = document.createElement('div');
  avatar.className = 'msg-avatar';
  avatar.textContent = '🤖';
  const dots = document.createElement('div');
  dots.className = 'typing-dots';
  dots.innerHTML = '<span></span><span></span><span></span>';
  div.appendChild(avatar);
  div.appendChild(dots);
  msgs.appendChild(div);
  msgs.scrollTop = msgs.scrollHeight;
  return div;
}

function removeTyping(el) { if (el && el.parentNode) el.parentNode.removeChild(el); }

// ── WELCOME ────────────────────────────────────────────────────────────────
function hideWelcome() {
  const w = document.getElementById('welcome');
  if (w) w.style.display = 'none';
}

// ── META UPDATE ────────────────────────────────────────────────────────────
function updateMeta(count) {
  document.getElementById('topbar-meta').textContent = `${count || state.msgCount} mensajes`;
}

// ── RUN COMMAND ────────────────────────────────────────────────────────────
async function runCmd(cmd, args = '') {
  if (state.loading) return;
  hideWelcome();
  state.loading = true;
  document.getElementById('send-btn').disabled = true;

  const displayText = `/${cmd}${args ? ' ' + args : ''}`;
  addMessage(displayText, 'user');
  const typing = addTyping();

  try {
    const r = await fetch('/api/comando', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({ session_id: state.sessionId, comando: cmd, args, modelo: state.model })
    });
    const d = await r.json();
    removeTyping(typing);
    addMessage(d.respuesta, 'bot');
    if (cmd === 'nuevo') { state.msgCount = 0; updateMeta(0); }
  } catch(e) {
    removeTyping(typing);
    addMessage('❌ Error de conexión.', 'bot');
  }

  state.loading = false;
  document.getElementById('send-btn').disabled = false;
}

// ── FOCUS CMD ──────────────────────────────────────────────────────────────
function focusCmd(cmd) {
  document.getElementById('cmd-input').value = '/' + cmd;
  document.getElementById('args-input').focus();
  document.getElementById('args-input').placeholder = `Tema para /${cmd}...`;
  closeSidebar();
}

// ── FILE UPLOAD ────────────────────────────────────────────────────────────
async function handleFileUpload(input) {
  const file = input.files[0];
  if (!file) return;

  hideWelcome();
  state.loading = true;
  document.getElementById('send-btn').disabled = true;
  addMessage(`📎 Subiendo: ${file.name}`, 'user');
  const typing = addTyping();

  const form = new FormData();
  form.append('archivo', file);
  form.append('session_id', state.sessionId);
  form.append('pregunta', document.getElementById('args-input').value.trim());
  form.append('modelo', state.model);

  try {
    const r = await fetch('/api/documento', { method: 'POST', body: form });
    const d = await r.json();
    removeTyping(typing);
    addMessage(d.respuesta || d.error, 'bot');
  } catch(e) {
    removeTyping(typing);
    addMessage('❌ Error al subir el archivo.', 'bot');
  }

  state.loading = false;
  document.getElementById('send-btn').disabled = false;
  input.value = '';
}

// ── IMAGE GENERATION ───────────────────────────────────────────────────────
async function generateImage(prompt) {
  hideWelcome();
  state.loading = true;
  document.getElementById('send-btn').disabled = true;
  addMessage(`🎨 /imagen ${prompt}`, 'user');
  const typing = addTyping();

  try {
    const r = await fetch('/api/imagen', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({ prompt })
    });
    const d = await r.json();
    removeTyping(typing);
    if (d.imagen) {
      const msgs = document.getElementById('messages');
      const div = document.createElement('div');
      div.className = 'msg bot';
      const avatar = document.createElement('div');
      avatar.className = 'msg-avatar';
      avatar.textContent = '🤖';
      const content = document.createElement('div');
      content.className = 'msg-content msg-image';
      const img = document.createElement('img');
      img.src = d.imagen;
      img.alt = d.prompt;
      img.style.borderRadius = '12px';
      img.style.maxWidth = '100%';
      img.style.maxHeight = '400px';
      const cap = document.createElement('div');
      cap.className = 'msg-time';
      cap.textContent = '🎨 ' + d.prompt;
      content.appendChild(img);
      content.appendChild(cap);
      div.appendChild(avatar);
      div.appendChild(content);
      msgs.appendChild(div);
      msgs.scrollTop = msgs.scrollHeight;
    } else {
      addMessage('❌ No se pudo generar la imagen. Intenta con otra descripción.', 'bot');
    }
  } catch(e) {
    removeTyping(typing);
    addMessage('❌ Error al generar imagen.', 'bot');
  }

  state.loading = false;
  document.getElementById('send-btn').disabled = false;
}

// Intercept /imagen command — integrated directly into sendMessage
async function generateImageFromCmd(prompt) {
  await generateImage(prompt);
}

// ── APA WEB MODAL ──────────────────────────────────────────────────────────
function openModal(id) { document.getElementById(id).classList.add('open'); document.getElementById('apa-url-input').focus(); }
function closeModal() { document.getElementById('apa-modal').classList.remove('open'); }

async function submitApaWeb() {
  const url = document.getElementById('apa-url-input').value.trim();
  if (!url.startsWith('http')) { showToast('Ingresa una URL válida', 'error'); return; }
  closeModal();
  hideWelcome();
  state.loading = true;
  document.getElementById('send-btn').disabled = true;
  addMessage(`/apa_web ${url}`, 'user');
  const typing = addTyping();

  try {
    const r = await fetch('/api/apa_web', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({ session_id: state.sessionId, url })
    });
    const d = await r.json();
    removeTyping(typing);
    addMessage(d.respuesta || d.error, 'bot');
  } catch(e) {
    removeTyping(typing);
    addMessage('❌ Error al procesar URL.', 'bot');
  }

  state.loading = false;
  document.getElementById('send-btn').disabled = false;
  document.getElementById('apa-url-input').value = '';
}

document.getElementById('apa-url-input').addEventListener('keydown', e => {
  if (e.key === 'Enter') submitApaWeb();
  if (e.key === 'Escape') closeModal();
});

// ── SIDEBAR TABS ───────────────────────────────────────────────────────────
function switchTab(tab, el) {
  document.querySelectorAll('.sidebar-tab').forEach(t => t.classList.remove('active'));
  document.querySelectorAll('.tab-content').forEach(t => t.classList.remove('active'));
  el.classList.add('active');
  document.getElementById(`tab-${tab}`).classList.add('active');
}

// ── SIDEBAR TOGGLE ─────────────────────────────────────────────────────────
function toggleSidebar() { document.getElementById('sidebar').classList.toggle('open'); }
function closeSidebar() {
  if (window.innerWidth <= 768) document.getElementById('sidebar').classList.remove('open');
}

// ── MODE & MODEL ───────────────────────────────────────────────────────────
function setMode(mode, btn) {
  state.mode = mode;
  document.querySelectorAll('.mode-btn').forEach(b => b.classList.remove('active'));
  btn.classList.add('active');
  const names = {libre:'💬 Chat Libre', academico:'🎓 Modo Académico', codigo:'💻 Modo Código', creativo:'🎨 Modo Creativo', tutor:'👨‍🏫 Modo Tutor'};
  document.getElementById('topbar-title').textContent = names[mode] || mode;
  showToast(`Modo ${mode} activado`, 'success');
}

function setModel(model) {
  state.model = model;
  showToast(`Motor: ${model.toUpperCase()}`, 'success');
}

// ── CONFIG SAVE ────────────────────────────────────────────────────────────
async function saveConfig() {
  const body = {
    groq_key: document.getElementById('groq-key').value,
    gemini_key: document.getElementById('gemini-key').value,
    openrouter_key: document.getElementById('or-key').value,
  };
  if (!body.groq_key && !body.gemini_key && !body.openrouter_key) {
    showToast('Ingresa al menos una clave API', 'error');
    return;
  }
  try {
    const r = await fetch('/api/config', {
      method: 'POST',
      headers: {'Content-Type': 'application/json'},
      body: JSON.stringify(body)
    });
    const d = await r.json();
    showToast(d.mensaje || 'Configuración guardada', 'success');
    ['groq-key','gemini-key','or-key'].forEach(id => document.getElementById(id).value = '');
    setTimeout(loadStatus, 1000);
  } catch(e) {
    showToast('Error al guardar configuración', 'error');
  }
}

// ── TOAST ──────────────────────────────────────────────────────────────────
let toastTimer;
function showToast(msg, type = '') {
  const el = document.getElementById('toast');
  el.textContent = msg;
  el.className = `show ${type}`;
  clearTimeout(toastTimer);
  toastTimer = setTimeout(() => { el.className = ''; }, 3000);
}

// ── CLOSE MODAL ON OVERLAY CLICK ───────────────────────────────────────────
document.getElementById('apa-modal').addEventListener('click', function(e) {
  if (e.target === this) closeModal();
});
</script>
</body>
</html>"""


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main() -> None:
    logger.info("=" * 64)
    logger.info("   %s Bot v%s — Iniciando servidor web", BOT_NAME, BOT_VERSION)
    logger.info("   Abre: http://localhost:%d", PORT)
    logger.info("=" * 64)

    if not GROQ_API_KEY and not GEMINI_API_KEY and not OPENROUTER_API_KEY:
        logger.warning(
            "⚠️  Ninguna clave API configurada. "
            "Usa el panel de Configuración en la interfaz web."
        )

    app.run(host="0.0.0.0", port=PORT, debug=False, threaded=True)


if __name__ == "__main__":
    main()
